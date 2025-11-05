#!/usr/bin/env python3
"""
Document Bot - Specialized Document Processing
Role: PDF, DOCX, TXT, CSV, XLSX document analysis using Gemini AI
"""
import os
import sys
import json
import time
import logging
import tempfile
from datetime import datetime
from typing import Dict, List, Optional

# Add parent directories to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(__file__))))

from dotenv import load_dotenv
load_dotenv()

from telegram import Bot
from shared.redis_utils import BotMessenger
from shared.gemini_client import GeminiAnalyzer
from shared.telegram_utils import TelegramClient

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('document_bot.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("document_bot")

# Configuration
DOCUMENT_BOT_TOKEN = os.getenv("DOCUMENT_BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
REDIS_HOST = os.getenv("REDIS_HOST", "localhost")
REDIS_PORT = os.getenv("REDIS_PORT", "6379")

# Initialize
messenger = BotMessenger("document_bot")
gemini = GeminiAnalyzer()
telegram_client = TelegramClient(DOCUMENT_BOT_TOKEN)


async def download_file_from_telegram(file_id: str, file_name: str) -> str:
    """Download file from Telegram and return local path"""
    try:
        # Get file from Telegram
        file = await telegram_client.get_file(file_id)

        # Create temp directory
        temp_dir = tempfile.mkdtemp()
        file_path = os.path.join(temp_dir, file_name)

        # Download file
        await file.download_to_drive(file_path)

        logger.info(f"Downloaded file: {file_name} to {file_path}")
        return file_path

    except Exception as e:
        logger.error(f"Error downloading file: {e}")
        raise


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file"""
    try:
        import PyPDF2

        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"

        logger.info(f"Extracted {len(text)} characters from PDF")
        return text

    except Exception as e:
        logger.error(f"Error extracting PDF text: {e}")
        raise


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file"""
    try:
        from docx import Document

        doc = Document(file_path)
        text = ""

        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"

        logger.info(f"Extracted {len(text)} characters from DOCX")
        return text

    except Exception as e:
        logger.error(f"Error extracting DOCX text: {e}")
        raise


def extract_text_from_txt(file_path: str) -> str:
    """Extract text from plain text file"""
    try:
        import chardet

        # Detect encoding
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'

        # Read with detected encoding
        with open(file_path, 'r', encoding=encoding) as file:
            text = file.read()

        logger.info(f"Extracted {len(text)} characters from TXT")
        return text

    except Exception as e:
        logger.error(f"Error extracting TXT text: {e}")
        raise


def extract_text_from_csv(file_path: str) -> str:
    """Extract text from CSV file"""
    try:
        import pandas as pd

        df = pd.read_csv(file_path)
        text = f"CSV 파일: {len(df)}행, {len(df.columns)}열\n\n"
        text += "컬럼: " + ", ".join(df.columns) + "\n\n"
        text += df.to_string()

        logger.info(f"Extracted {len(text)} characters from CSV")
        return text

    except Exception as e:
        logger.error(f"Error extracting CSV text: {e}")
        raise


def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from XLSX file"""
    try:
        import pandas as pd

        # Try to read all sheets
        excel_file = pd.ExcelFile(file_path)
        text = f"XLSX 파일: {len(excel_file.sheet_names)}개 시트\n\n"

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            text += f"시트: {sheet_name} ({len(df)}행, {len(df.columns)}열)\n"
            text += "컬럼: " + ", ".join(df.columns) + "\n\n"

        logger.info(f"Extracted {len(text)} characters from XLSX")
        return text

    except Exception as e:
        logger.error(f"Error extracting XLSX text: {e}")
        raise


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text = f"PPTX 파일: {len(prs.slides)}개 슬라이드\n\n"

        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"슬라이드 {slide_num}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"

        logger.info(f"Extracted {len(text)} characters from PPTX")
        return text

    except Exception as e:
        logger.error(f"Error extracting PPTX text: {e}")
        raise


async def process_document_task(task_data: Dict):
    """Process document analysis task"""
    try:
        chat_id = task_data.get('chat_id')
        file_data = task_data.get('file_data', {})
        file_id = file_data.get('file_id')
        file_name = file_data.get('file_name', 'document')

        logger.info(f"Processing document: {file_name} for chat {chat_id}")

        # Send progress update
        messenger.notify_progress(chat_id, "문서를 다운로드하는 중...")

        # Download file
        file_path = await download_file_from_telegram(file_id, file_name)

        # Extract text based on file type
        messenger.notify_progress(chat_id, "텍스트를 추출하는 중...")

        file_ext = os.path.splitext(file_name)[1].lower()
        extracted_text = ""

        if file_ext == '.pdf':
            extracted_text = extract_text_from_pdf(file_path)
        elif file_ext == '.docx':
            extracted_text = extract_text_from_docx(file_path)
        elif file_ext == '.txt':
            extracted_text = extract_text_from_txt(file_path)
        elif file_ext == '.csv':
            extracted_text = extract_text_from_csv(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            extracted_text = extract_text_from_xlsx(file_path)
        elif file_ext == '.pptx':
            extracted_text = extract_text_from_pptx(file_path)
        else:
            # Try as plain text
            extracted_text = extract_text_from_txt(file_path)

        # Limit text length for AI processing
        max_length = 10000
        if len(extracted_text) > max_length:
            extracted_text = extracted_text[:max_length] + "\n\n[텍스트가 길어서 일부만 분석했습니다]"

        # Send progress update
        messenger.notify_progress(chat_id, "Gemini AI로 분석하는 중...")

        # Analyze with Gemini AI
        summary = gemini.analyze_document(extracted_text)

        # Prepare result
        result = {
            "text": extracted_text,
            "summary": summary,
            "file_name": file_name,
            "processed_at": datetime.now().isoformat()
        }

        # Send result to main bot
        messenger.send_result(chat_id, result)

        # Clean up
        try:
            os.remove(file_path)
            os.rmdir(os.path.dirname(file_path))
        except:
            pass

        logger.info(f"Completed document analysis for chat {chat_id}")

    except Exception as e:
        logger.error(f"Error processing document task: {e}")
        # Send error result
        error_result = {
            "error": str(e),
            "file_name": task_data.get('file_data', {}).get('file_name', 'unknown')
        }
        messenger.send_result(task_data.get('chat_id'), error_result)


async def listen_for_tasks():
    """Listen for document processing tasks"""
    logger.info("Document bot started, listening for tasks...")

    pubsub = messenger.pubsub
    pubsub.subscribe("document_tasks")

    for message in pubsub.listen():
        if message['type'] == 'message':
            try:
                data = json.loads(message['data'])
                await process_document_task(data)
            except Exception as e:
                logger.error(f"Error processing task: {e}")


async def main():
    """Main function"""
    print("=== Document Bot (Document Processing) ===")

    if not DOCUMENT_BOT_TOKEN:
        print("❌ ERROR: DOCUMENT_BOT_TOKEN is missing")
        print("Please set DOCUMENT_BOT_TOKEN in .env file")
        return

    if not GEMINI_API_KEY:
        print("⚠️ WARNING: GEMINI_API_KEY is missing - AI features will be disabled")

    try:
        # Test Telegram connection
        bot = Bot(token=MAIN_BOT_TOKEN)
        await bot.get_me()
        print("✅ Telegram connection successful")
    except Exception as e:
        print(f"❌ ERROR: Failed to connect to Telegram: {e}")
        return

    try:
        # Start listening for tasks
        await listen_for_tasks()
    except KeyboardInterrupt:
        print("\n👋 Shutting down...")
    finally:
        messenger.close()


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
