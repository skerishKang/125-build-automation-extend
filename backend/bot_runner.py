#!/usr/bin/env python3
"""
125 Build Automation - Telegram Bot Runner (Gemini 2.0 Flash Multimodal)
- Single file handling text/document/image/voice with Gemini 2.0 Flash
- Free chat with memory (Supabase optional)
- Document/Image/Voice processed directly with Gemini's multimodal capabilities
- Google Drive bidirectional sync
"""
import os
import sys
import logging
from datetime import datetime
from typing import Dict, List, Any
import tempfile
import asyncio

from dotenv import load_dotenv
load_dotenv(os.path.join(os.path.dirname(__file__), ".env"))

TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_ANON_KEY")

# logging
os.makedirs("logs", exist_ok=True)
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(os.path.join("logs", "bot_runner.log")),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("unified_bot")

# Disable httpx logging to prevent token exposure
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("telegram").setLevel(logging.WARNING)

# telegram
try:
    from telegram import Update
    from telegram.constants import ChatAction
    from telegram.ext import Application, CommandHandler, MessageHandler, ContextTypes, filters
except ImportError:
    logger.error("python-telegram-bot is not installed. pip install python-telegram-bot==21.6")
    sys.exit(1)

# gemini (multimodal)
gemini_model = None
if GEMINI_API_KEY:
    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        gemini_model = genai.GenerativeModel('gemini-2.5-flash')
        logger.info("Using Gemini 2.5 Flash (multimodal)")
    except Exception as e:
        logger.error(f"Gemini setup failed: {e}")
else:
    logger.warning("GEMINI_API_KEY not set; chat will be disabled")

# supabase (optional memory)
supabase = None
if SUPABASE_URL and SUPABASE_KEY:
    try:
        from supabase import create_client
        supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
    except Exception as e:
        logger.warning(f"Supabase init failed: {e}")

# in-memory recent docs (fallback)
recent_documents: Dict[int, List[Dict[str, Any]]] = {}

# User patterns and learning system
user_patterns: Dict[str, Dict[str, Any]] = {}
user_modes: Dict[str, str] = {}  # 'beginner', 'advanced', 'hybrid'

# Multi-turn conversation context
conversation_contexts: Dict[str, Dict[str, Any]] = {}  # user_id -> context

# Load user patterns from file
def load_user_patterns():
    """Load user patterns from file"""
    global user_patterns
    try:
        if os.path.exists(USER_PATTERN_FILE):
            with open(USER_PATTERN_FILE, 'r') as f:
                user_patterns = json.load(f)
    except Exception as e:
        logger.warning(f"Failed to load user patterns: {e}")

# Save user patterns to file
def save_user_patterns():
    """Save user patterns to file"""
    try:
        with open(USER_PATTERN_FILE, 'w') as f:
            json.dump(user_patterns, f, indent=2)
    except Exception as e:
        logger.error(f"Failed to save user patterns: {e}")

# Load user modes from file
def load_user_modes():
    """Load user modes from file"""
    global user_modes
    try:
        if os.path.exists(USER_MODE_FILE):
            with open(USER_MODE_FILE, 'r') as f:
                user_modes = json.load(f)
    except Exception as e:
        logger.warning(f"Failed to load user modes: {e}")

# Save user modes to file
def save_user_modes():
    """Save user modes to file"""
    try:
        with open(USER_MODE_FILE, 'w') as f:
            json.dump(user_modes, f, indent=2)
    except Exception as e:
        logger.error(f"Failed to save user modes: {e}")

# Load conversation contexts from file
def load_conversation_contexts():
    """Load conversation contexts from file"""
    global conversation_contexts
    try:
        if os.path.exists(CONVERSATION_CONTEXT_FILE):
            with open(CONVERSATION_CONTEXT_FILE, 'r') as f:
                conversation_contexts = json.load(f)
    except Exception as e:
        logger.warning(f"Failed to load conversation contexts: {e}")

# Save conversation contexts to file
def save_conversation_contexts():
    """Save conversation contexts to file"""
    try:
        with open(CONVERSATION_CONTEXT_FILE, 'w') as f:
            json.dump(conversation_contexts, f, indent=2)
    except Exception as e:
        logger.error(f"Failed to save conversation contexts: {e}")

# Initialize all user data on startup
load_user_patterns()
load_user_modes()
load_conversation_contexts()


# ========== Intent Analysis Engine ==========

async def analyze_user_intent(message: str, user_id: str) -> Dict[str, Any]:
    """Analyze user intent using Gemini AI"""
    if not gemini_model:
        return {
            "action": "chat",
            "confidence": 0.0,
            "parameters": {},
            "natural_response": "AI 분석 엔진이 비활성화되어 있어요."
        }

    try:
        # Get user patterns
        patterns = user_patterns.get(user_id, {})
        
        # Build analysis prompt
        analysis_prompt = f"""
사용자 메시지: "{message}"
사용자 패턴: {json.dumps(patterns, ensure_ascii=False)}

다음 JSON 형태로 의도를 분석해주세요:
{{
    "action": "실행할 액션명",
    "parameters": {{"매개변수": "값"}},
    "confidence": 0.95,
    "natural_response": "자연어로 반응할 메시지",
    "requires_confirmation": false,
    "follow_up_question": "추가 질문이 필요한 경우"
}}

지원하는 액션들:
- gmail_list: Gmail 메일 목록 보기 (parameters: {{"max_results": 10}})
- gmail_on: Gmail 모니터링 시작
- gmail_off: Gmail 모니터링 중지
- cal_today: 오늘 일정 보기
- cal_tomorrow: 내일 일정 보기
- cal_week: 이번 주 일정 보기
- cal_search: 일정 검색 (parameters: {{"query": "검색어"}})
- cal_on: Calendar 모니터링 시작
- cal_off: Calendar 모니터링 중지
- drive_list: 드라이브 파일 목록
- excel_extract: 엑셀 데이터 추출 (parameters: {{"sheet": "시트명"}})
- summarize: 문서 요약
- reply: 메일 답장 작성 (parameters: {{"email_id": "id"}})
- chat: 일반 대화
- excel_create: 엑셀 파일 생성

패턴 인식 규칙:
- "메일" + "확인/보기/체크" = gmail_list
- "메일" + "답장/회신" = reply
- "일정" + "오늘/오늘일정" = cal_today
- "일정" + "내일/내일일정" = cal_tomorrow
- "일정" + "이번주/주간" = cal_week
- "일정" + "검색/찾기" = cal_search
- "드라이브" + "목록/파일" = drive_list
- "엑셀" + "요약/분석" = summarize
- "답장" + "써줘/작성" = reply

항상 한국어로만 답변하고, confidence는 0.0-1.0 사이의 실수로 설정해주세요.
"""

        response = gemini_model.generate_content(analysis_prompt)
        
        # Extract JSON from response
        import re
        json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
        if json_match:
            result = json.loads(json_match.group())
            
            # Learn user pattern if confidence > 0.8
            if result.get('confidence', 0) > 0.8 and result.get('action') != 'chat':
                learn_user_pattern(user_id, message, result['action'])
            
            return result
        else:
            logger.warning(f"Failed to extract JSON from Gemini response: {response.text}")
            return {
                "action": "chat",
                "confidence": 0.0,
                "parameters": {},
                "natural_response": "무슨 뜻인지 잘 이해하지 못했어요. 명령어를 사용해주세요!"
            }

    except Exception as e:
        logger.error(f"Intent analysis error: {e}")
        return {
            "action": "chat",
            "confidence": 0.0,
            "parameters": {},
            "natural_response": "AI 분석 중 오류가 발생했어요. 다시 시도해주세요."
        }


def learn_user_pattern(user_id: str, natural_text: str, action_taken: str):
    """Learn user pattern for better intent recognition"""
    global user_patterns
    
    if user_id not in user_patterns:
        user_patterns[user_id] = {
            "gmail_patterns": [],
            "calendar_patterns": [],
            "drive_patterns": [],
            "excel_patterns": [],
            "chat_patterns": [],
            "learned_count": 0
        }
    
    # Determine pattern category based on action
    category = "chat_patterns"
    if action_taken.startswith("gmail_"):
        category = "gmail_patterns"
    elif action_taken.startswith("cal_"):
        category = "calendar_patterns"
    elif action_taken.startswith("drive"):
        category = "drive_patterns"
    elif action_taken.startswith("excel_"):
        category = "excel_patterns"
    
    # Add pattern if not already exists
    patterns = user_patterns[user_id][category]
    if natural_text not in patterns:
        patterns.append(natural_text)
        user_patterns[user_id]["learned_count"] = user_patterns[user_id].get("learned_count", 0) + 1
        
        # Save to file
        save_user_patterns()
        
        logger.info(f"Learned pattern for user {user_id}: {natural_text} -> {action_taken}")


async def execute_action(update: Update, intent: Dict[str, Any], context: ContextTypes.DEFAULT_TYPE):
    """Execute action based on intent"""
    action = intent.get('action')
    parameters = intent.get('parameters', {})
    
    try:
        if action == 'gmail_list':
            await handle_gmail_list(update, context)
        elif action == 'gmail_on':
            await handle_gmail_on(update, context)
        elif action == 'gmail_off':
            await handle_gmail_off(update, context)
        elif action == 'cal_today':
            await handle_cal_today(update, context)
        elif action == 'cal_tomorrow':
            await handle_cal_tomorrow(update, context)
        elif action == 'cal_week':
            await handle_cal_week(update, context)
        elif action == 'cal_search':
            query = parameters.get('query', '')
            if query:
                # Simulate command with arguments
                context.args = query.split()
                await handle_cal_search(update, context)
        elif action == 'cal_on':
            await handle_cal_on(update, context)
        elif action == 'cal_off':
            await handle_cal_off(update, context)
        elif action == 'drive_list':
            await handle_drive_list(update, context)
        elif action == 'reply':
            # For now, just acknowledge
            await reply_text(update, "✏️ 메일 답장 기능을 실행할게요. 어떤 메일에 답장할지 알려주세요!")
        elif action == 'summarize':
            await reply_text(update, "📄 문서 요약 기능을 실행할게요! 파일을 업로드해주세요!")
        else:
            await handle_text(update, context)  # Fallback to general chat

    except Exception as e:
        logger.error(f"Action execution error: {e}")
        await reply_text(update, f"실행 중 오류가 발생했어요: {str(e)[:100]}")


async def confirm_and_execute(update: Update, intent: Dict[str, Any], context: ContextTypes.DEFAULT_TYPE):
    """Confirm with user before executing ambiguous intent"""
    natural_response = intent.get('natural_response', '이렇게 이해하면 될까요?')
    action = intent.get('action')
    confidence = intent.get('confidence', 0.5)
    
    # Build confirmation message
    action_names = {
        'gmail_list': '📧 Gmail 메일 확인하기',
        'cal_today': '📅 오늘 일정 확인하기',
        'cal_tomorrow': '📅 내일 일정 확인하기',
        'cal_week': '📅 이번 주 일정 확인하기',
        'drive_list': '📁 드라이브 파일 목록 보기'
    }
    
    action_name = action_names.get(action, action)
    
    confirm_text = f"""
🤔 **이렇게 이해해도 될까요?**

{natural_response}

→ **{action_name}** 실행

💡 아니면 다른 뜻이었나요?
    """.strip()
    
    # Store context for follow-up
    user_id = str(update.effective_user.id)
    conversation_contexts[user_id] = {
        'pending_intent': intent,
        'timestamp': datetime.now().isoformat()
    }
    save_conversation_contexts()
    
    # Add inline keyboard with Yes/No buttons
    from telegram import InlineKeyboardButton, InlineKeyboardMarkup
    
    keyboard = [
        [
            InlineKeyboardButton("✅ 맞음", callback_data="intent_yes"),
            InlineKeyboardButton("❌ 아님", callback_data="intent_no")
        ]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)
    
    await context.bot.send_message(
        chat_id=update.effective_chat.id,
        text=confirm_text,
        reply_markup=reply_markup
    )


async def suggest_smart_actions(file_type: str, file_name: str) -> str:
    """Suggest smart actions based on file type"""
    if file_type in ['.xlsx', '.xls']:
        suggestions = [
            "💡 데이터 요약 생성하기",
            "💡 차트 자동 만들기",
            "💡 특정 데이터 추출하기",
            "💡 다른 형태로 변환하기"
        ]
    elif file_type == '.pdf':
        suggestions = [
            "💡 문서 요약하기",
            "💡 주요 내용 추출하기",
            "💡 번역하기",
            "💡 Word로 변환하기"
        ]
    elif file_type in ['.txt', '.md', '.py', '.js']:
        suggestions = [
            "💡 코드 분석하기",
            "💡 최적화 제안 받기",
            "💡 문서화 자동 생성"
        ]
    elif file_type in ['.png', '.jpg', '.jpeg']:
        suggestions = [
            "💡 이미지 설명받기",
            "💡 텍스트 추출하기 (OCR)",
            "💡 이미지 개선하기"
        ]
    else:
        suggestions = [
            "💡 내용 요약하기",
            "💡 형식 분석하기"
        ]
    
    return f"📁 **{file_name}**을 받았어요!\n\n" + "\n".join(suggestions)


# ========== Enhanced Multi-turn Conversation ==========

async def handle_follow_up_response(update: Update, context: ContextTypes.DEFAULT_TYPE, user_id: str):
    """Handle follow-up responses in multi-turn conversation"""
    if user_id not in conversation_contexts:
        await reply_text(update, "대화 맥락을 찾을 수 없어요. 다시 시작해주세요!")
        return
    
    context_data = conversation_contexts[user_id]
    pending_intent = context_data.get('pending_intent')
    
    if not pending_intent:
        await reply_text(update, "대화 맥락을 찾을 수 없어요. 다시 시작해주세요!")
        return
    
    # Execute the pending action
    await execute_action(update, pending_intent, context)
    
    # Clear context
    del conversation_contexts[user_id]
    save_conversation_contexts()


# Inline keyboard callback handler (add to main handler)
async def handle_callback_query(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle inline keyboard callbacks"""
    query = update.callback_query
    await query.answer()

    user_id = str(query.from_user.id)
    data = query.data

    # Handle intent confirmation callbacks
    if data == "intent_yes":
        if user_id in conversation_contexts and 'pending_intent' in conversation_contexts[user_id]:
            pending_intent = conversation_contexts[user_id]['pending_intent']
            await execute_action(update, pending_intent, context)
            del conversation_contexts[user_id]
            save_conversation_contexts()
    elif data == "intent_no":
        await query.edit_message_text("😔 아, 실수했네요! 다시 말해주세요!")
        if user_id in conversation_contexts:
            del conversation_contexts[user_id]
            save_conversation_contexts()

    # Handle reply-related callbacks
    elif data.startswith("send_reply_") or data.startswith("edit_reply_") or data.startswith("regenerate_reply_") or data == "cancel_reply":
        await handle_reply_callback(update, context)

# Smart audio processing configuration
SHORT_AUDIO_THRESHOLD = int(os.getenv("SHORT_AUDIO_THRESHOLD", "30"))  # 30초 이하
LONG_AUDIO_THRESHOLD = int(os.getenv("LONG_AUDIO_THRESHOLD", "300"))  # 5분 이상
MID_LENGTH_MODEL = os.getenv("MID_LENGTH_AUDIO", "gemini")  # 30초-5분 기본

# Drive monitoring configuration
DRIVE_MONITOR_INTERVAL = int(os.getenv("DRIVE_MONITOR_INTERVAL", "300"))  # 5분 (300초)
ENABLE_DRIVE_MONITORING = os.getenv("ENABLE_DRIVE_MONITORING", "true").lower() == "true"

# Gmail monitoring configuration
GMAIL_MONITOR_INTERVAL = int(os.getenv("GMAIL_MONITOR_INTERVAL", "300"))  # 기본 5분 간격

# Usage telemetry counters
usage_metrics = {
    "text_messages": 0,
    "documents_saved": 0,
    "photos_processed": 0,
    "voices_processed": 0,
    "drive_sync_runs": 0,
    "drive_downloads": 0,
    "gmail_notifications": 0,
    "calendar_alerts_sent": 0
}

# Global application instance for Drive monitoring
_app_instance = None

# Drive monitoring state control
drive_monitoring_state = {
    "enabled": False,
    "thread": None,
    "last_check": None,
    "total_files": 0,
    "start_time": None,
    "interval": DRIVE_MONITOR_INTERVAL,
    "recent_changes": deque(maxlen=20),
    "last_manual_sync": None
}

# Gmail monitoring state control
gmail_monitoring_state = {
    "enabled": False,
    "thread": None,
    "last_check": None,
    "total_emails": 0,
    "start_time": None,
    "interval": GMAIL_MONITOR_INTERVAL,
    "auth_mode": "unknown",
    "recent_emails": deque(maxlen=50)
}

# Calendar monitoring state control
calendar_monitoring_state = {
    "enabled": False,
    "thread": None,
    "last_check": None,
    "total_alerts": 0,
    "start_time": None,
    "alerted_events": set()  # Track alerted event IDs
}


def get_audio_duration(ogg_path: str) -> float:
    """Get audio duration in seconds using ffprobe (if available) or estimate"""
    try:
        import subprocess
        result = subprocess.run(
            ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
             "-of", "csv=p=0", ogg_path],
            capture_output=True, text=True, check=True
        )
        return float(result.stdout.strip())
    except Exception:
        # Fallback: estimate based on file size (rough)
        # ~1MB per minute at 64kbps
        size_mb = os.path.getsize(ogg_path) / (1024 * 1024)
        return size_mb * 60 * 0.7  # conservative estimate


def format_plain(text: str, max_len: int = 1200) -> str:
    """Format Gemini response to Telegram-friendly plain text"""
    import re
    # Remove code blocks
    text = re.sub(r"```.*?```", "", text, flags=re.DOTALL)
    # Remove tables
    text = re.sub(r"\|.*\|", "", text)
    # Remove header symbols (keep line breaks)
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    # List symbols (keep line breaks)
    text = re.sub(r"^\s*[-*•]\s*", "• ", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+\.\s*", "• ", text, flags=re.MULTILINE)
    # Remove bold/italic
    text = text.replace("**", "").replace("*", "")
    # Remove backticks
    text = text.replace("`", "'")
    # Clean up multiple line breaks (max 2)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip trailing spaces
    text = "\n".join(line.rstrip() for line in text.split("\n"))
    # Strip leading/trailing spaces
    text = text.strip()
    # Length limit with ...
    if len(text) > max_len:
        text = text[:max_len] + "…"
    return text


# Thread pool for CPU-intensive operations (transcription, etc.)
from concurrent.futures import ThreadPoolExecutor
audio_executor = ThreadPoolExecutor(max_workers=2, thread_name_prefix="audio_processing")


async def _action_indicator(context: ContextTypes.DEFAULT_TYPE, chat_id: int, action: str, stop_event: asyncio.Event):
    try:
        while not stop_event.is_set():
            try:
                await context.bot.send_chat_action(chat_id=chat_id, action=action)
            except Exception:
                pass
            # Telegram은 5초 동안 액션 유지. 4초 주기로 새로 송신.
            try:
                await asyncio.wait_for(stop_event.wait(), timeout=4.0)
            except asyncio.TimeoutError:
                continue
    except Exception:
        pass

class ActionIndicator:
    def __init__(self, context: ContextTypes.DEFAULT_TYPE, chat_id: int, action: str):
        self.context = context
        self.chat_id = chat_id
        self.action = action
        self.stop_event = asyncio.Event()
        self.task: asyncio.Task | None = None

    async def __aenter__(self):
        self.task = asyncio.create_task(_action_indicator(self.context, self.chat_id, self.action, self.stop_event))
        return self

    async def __aexit__(self, exc_type, exc, tb):
        self.stop_event.set()
        if self.task:
            try:
                await asyncio.wait_for(self.task, timeout=1.5)
            except Exception:
                self.task.cancel()


async def save_memory(user_id: str, username: str, message: str, response: str):
    if not supabase:
        return
    try:
        supabase.table("conversations").insert({
            "user_id": user_id,
            "username": username,
            "message": message,
            "response": response,
            "created_at": datetime.utcnow().isoformat()
        }).execute()
    except Exception as e:
        logger.warning(f"save_memory failed: {e}")


async def fetch_memory(user_id: str, limit: int = 8) -> List[Dict[str, str]]:
    if not supabase:
        return []
    try:
        res = supabase.table("conversations").select("message,response,created_at").eq("user_id", user_id).order("created_at", desc=True).limit(limit).execute()
        return list(reversed(res.data or []))
    except Exception as e:
        logger.warning(f"fetch_memory failed: {e}")
        return []


async def reply_text(update: Update, text: str):
    # Prevent telegram 409: retry with slight delay on 409
    try:
        await update.message.reply_text(text)
    except Exception as e:
        logger.warning(f"reply_text failed: {e}")
        await asyncio.sleep(0.8)
        try:
            await update.message.reply_text(text[:4000])
        except Exception:
            pass


async def handle_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    name = update.effective_user.first_name or "사용자"
    monitoring_status = "🔄 Drive 자동 모니터링" if ENABLE_DRIVE_MONITORING else "📋 Manual Drive 체크"
    await reply_text(update,
        f"안녕하세요 {name}님! 👋\n\n"
        "이 봇은 Gemini 2.5 Flash 기반 \"올인원\"입니다.\n"
        "- 자유 대화 (메모리 포함)\n"
        "- 문서/이미지/음성 멀티모달 처리\n"
        "- Google Drive 양방향 동기화\n"
        "- Gmail 실시간 감시 및 AI 요약\n"
        f"- {monitoring_status}\n\n"
        "📂 **Drive 명령어**: /drive\n"
        "📧 **Gmail 명령어**: /gmail_on, /gmail_off")


async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = (update.message.text or "").strip()
    if not text or text.startswith('/'):
        return

    user_id = str(update.effective_user.id)
    username = update.effective_user.first_name or "사용자"

    # Check if this is a reply editing mode
    if user_id in conversation_contexts and conversation_contexts[user_id].get('editing_reply'):
        # Handle reply editing
        await handle_reply_edit(update, context, text)
        return

    # Check if this is a follow-up response
    if user_id in conversation_contexts:
        await handle_follow_up_response(update, context, user_id)
        return

    # Get user mode (hybrid by default)
    user_mode = user_modes.get(user_id, 'hybrid')

    # If mode is advanced, only process with Gemini (no intent analysis)
    if user_mode == 'advanced':
        await handle_text_advanced(update, context, text, user_id, username)
        return

    # For hybrid and beginner modes, use intent analysis
    # Analyze user intent
    intent = await analyze_user_intent(text, user_id)

    confidence = intent.get('confidence', 0.0)
    action = intent.get('action', 'chat')

    # If high confidence and it's not just chat, execute directly
    if confidence > 0.8 and action != 'chat':
        natural_response = intent.get('natural_response', f'{action} 실행할게요!')
        await reply_text(update, natural_response)
        await execute_action(update, intent, context)

        # Save to memory
        await save_memory(user_id, username, text, natural_response)

    # If medium confidence, ask for confirmation
    elif confidence > 0.5:
        await confirm_and_execute(update, intent, context)

    # Low confidence or just chat, use Gemini
    else:
        await handle_text_advanced(update, context, text, user_id, username)


async def handle_reply_edit(update: Update, context: ContextTypes.DEFAULT_TYPE, new_text: str):
    """Handle reply editing mode"""
    user_id = str(update.effective_user.id)

    if user_id not in conversation_contexts:
        await reply_text(update, "❌ 대화 맥락을 찾을 수 없습니다.")
        return

    # Get pending reply
    if 'pending_reply' not in conversation_contexts[user_id]:
        await reply_text(update, "❌ 답장 정보를 찾을 수 없습니다.")
        return

    # Update reply draft
    conversation_contexts[user_id]['pending_reply']['draft'] = new_text
    conversation_contexts[user_id]['editing_reply'] = False
    save_conversation_contexts()

    # Show updated reply with buttons
    reply_data = conversation_contexts[user_id]['pending_reply']

    final_message = f"""
✏️ **수정된 답장**:

{reply_data['draft']}

어떻게 하시겠어요?
    """.strip()

    from telegram import InlineKeyboardButton, InlineKeyboardMarkup

    keyboard = [
        [
            InlineKeyboardButton("📤 바로 보내기", callback_data=f"send_reply_{reply_data['email_id']}"),
            InlineKeyboardButton("✏️ 다시 수정", callback_data=f"edit_reply_{reply_data['email_id']}")
        ],
        [
            InlineKeyboardButton("🔄 다른 톤으로", callback_data=f"regenerate_reply_{reply_data['email_id']}"),
            InlineKeyboardButton("❌ 취소", callback_data="cancel_reply")
        ]
    ]
    reply_markup = InlineKeyboardMarkup(keyboard)

    await reply_text(update, final_message, reply_markup=reply_markup)


async def handle_text_advanced(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str, user_id: str, username: str):
    """Handle text with Gemini AI (for advanced mode or low confidence intent)"""
    if not GEMINI_API_KEY or not gemini_model:
        await reply_text(update, "Gemini 설정이 없어 대화가 비활성화되어 있어요.")
        return

    # Fetch memory and build context
    memory = await fetch_memory(user_id)
    context_lines = []
    if memory:
        context_lines.append("[이전 대화 맥락]")
        for m in memory:
            context_lines.append(f"User: {m['message']}")
            context_lines.append(f"Assistant: {m['response']}")
        context_lines.append("")

    # Smart keyword detection for response length
    short_keywords = ["요약", "간단히", "짧게", "요약", "간단"]
    long_keywords = ["자세히", "구체적으로", "설명", "상세히", "자세한"]
    is_short_question = any(keyword in text for keyword in short_keywords)
    is_long_question = any(keyword in text for keyword in long_keywords)

    # Smart prompt
    if is_long_question:
        prompt_style = "자세하고 구체적으로 설명해 주세요."
    elif is_short_question:
        prompt_style = "간단히 요약해 주세요."
    else:
        prompt_style = "간단히 요약해 주세요. 더 자세히 필요하면 추가 요청해 주세요."

    prompt = "\n".join(context_lines + [
        f"현재 사용자 메시지: {text}",
        f"답변 스타일: {prompt_style}",
        "항상 한국어로만 답변하고, Markdown 표/코드블록 없이 간결한 문장으로 답하세요."
    ])

    # Cumulative progress messages
    progress_messages = []
    progress_messages.append(await update.message.reply_text("💬 답변 생성 중… [10%]"))

    indicator = ActionIndicator(context, update.effective_chat.id, ChatAction.TYPING)
    await indicator.__aenter__()

    progress_messages.append(await update.message.reply_text("🧠 Gemini 2.5 Flash 분석 중… [50%]"))

    try:
        # 2) 블로킹 추론을 스레드로 오프로딩하여 동시 메시지 처리 유지
        def _call_gemini():
            resp = gemini_model.generate_content(prompt)
            return resp.text.strip()
        raw = await asyncio.to_thread(_call_gemini)
        answer = format_plain(raw)
        logger.info(f"Bot replied ({len(answer)} chars): {answer[:100]}...")
    except Exception as e:
        logger.error(f"Gemini error: {e}")
        answer = "죄송해요, 지금은 답변을 생성할 수 없어요."
    finally:
        await indicator.__aexit__(None, None, None)

    progress_messages.append(await update.message.reply_text("✅ 답변 완성! [100%]"))

    # 4) Send final result as new message
    final_text = f"{answer}"
    await reply_text(update, final_text)

    await save_memory(user_id, username, text, answer)


async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not GEMINI_API_KEY or not gemini_model:
        await reply_text(update, "Gemini 설정이 없어 이미지 분석이 비활성화되어 있어요.")
        return

    # Cumulative progress messages
    progress_messages = []
    progress_messages.append(await update.message.reply_text("📷 이미지를 받았어요. 분석 중… [0%]"))

    try:
        photo = update.message.photo[-1]
        file = await context.bot.get_file(photo.file_id)
        tmp = os.path.join(tempfile.gettempdir(), f"{photo.file_id}.jpg")
        photo_indicator = ActionIndicator(context, update.effective_chat.id, ChatAction.UPLOAD_PHOTO)
        await photo_indicator.__aenter__()
        await file.download_to_drive(tmp)

        # Step update: download complete
        progress_messages.append(await update.message.reply_text("📷 이미지 다운로드 완료. 멀티모달 분석 중… [50%]"))

        # Use Gemini's multimodal capability - upload image directly
        import google.generativeai as genai
        image_part = {"mime_type": "image/jpeg", "data": open(tmp, "rb").read()}

        prompt = "다음 이미지를 한국어로 설명하는 캡션을 작성해줘. 이미지의 주요 내용, 색감/분위기, 맥락을 간결하게 설명해주세요."
        prompt += "\n\n항상 한국어로만 답변하고, Markdown 표/코드블록 없이 간결한 문장으로 답하세요."

        # Multimodal call with image
        response = gemini_model.generate_content([prompt, image_part])
        answer = response.text.strip()
        answer = format_plain(answer)

        progress_messages.append(await update.message.reply_text("✅ 이미지 분석 완료! [100%]"))

        final_text = f"🖼️ 이미지 설명:\n{answer}"
        await reply_text(update, final_text)
    except Exception as e:
        logger.error(f"photo error: {e}")
        await reply_text(update, "이미지 처리에 실패했어요.")
    finally:
        # Clean up temp file
        try:
            if 'tmp' in locals():
                os.remove(tmp)
        except Exception:
            pass
        try:
            if 'photo_indicator' in locals():
                await photo_indicator.__aexit__(None, None, None)
        except Exception:
            pass


async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not GEMINI_API_KEY or not gemini_model:
        await reply_text(update, "Gemini 설정이 없어 음성 처리가 비활성화되어 있어요.")
        return

    # Immediate acknowledgment + background processing message
    ack_msg = None
    try:
        ack_msg = await update.message.reply_text(
            "🎤 음성을 받았어요. 백그라운드에서 처리 중입니다! "
            "다른 메시지도 바로 보낼 수 있어요. 😊"
        )
    except Exception:
        ack_msg = None

    chat_id = update.effective_chat.id
    user_id = str(update.effective_user.id)
    username = update.effective_user.first_name or "사용자"

    # Create background task for voice processing (non-blocking)
    asyncio.create_task(process_voice_background(update, context, chat_id, user_id, username, ack_msg))


async def process_voice_background(update, context, chat_id, user_id, username, ack_msg):
    """Process voice in background - non-blocking, allows immediate responses"""
    voice = update.message.voice
    file = await context.bot.get_file(voice.file_id)
    ogg_path = os.path.join(tempfile.gettempdir(), f"{voice.file_id}.ogg")
    wav_path = os.path.join(tempfile.gettempdir(), f"{voice.file_id}.wav")

    # Progress tracking for voice processing
    progress_messages = []

    try:
        # Download voice file
        await file.download_to_drive(ogg_path)
        progress_messages.append(await context.bot.send_message(chat_id, "📥 음성 파일 다운로드 완료. [20%]"))

        # Get audio duration
        duration = get_audio_duration(ogg_path)
        progress_messages.append(await context.bot.send_message(chat_id, f"⏱️ 음성 길이 분석: {duration:.1f}초. 처리 방식 결정 중... [40%]"))

        # Select model based on duration
        if duration <= SHORT_AUDIO_THRESHOLD:
            # SHORT: Use Gemini 2.5 Flash (multimodal, fast)
            result = await process_with_gemini_multimodal(ogg_path, duration, chat_id, context, progress_messages)
            mode = "Gemini 2.5 Flash (멀티모달)"
        elif duration >= LONG_AUDIO_THRESHOLD:
            # LONG: Use Whisper + Gemini (accurate, free)
            result = await process_with_whisper_gemini(ogg_path, wav_path, duration, chat_id, context, progress_messages)
            mode = "Whisper + Gemini (정확도 최적화)"
        else:
            # MID: Use environment setting
            if MID_LENGTH_MODEL == "gemini":
                result = await process_with_gemini_multimodal(ogg_path, duration, chat_id, context, progress_messages)
                mode = "Gemini 2.5 Flash (멀티모달)"
            else:
                result = await process_with_whisper_gemini(ogg_path, wav_path, duration, chat_id, context, progress_messages)
                mode = "Whisper + Gemini (정확도 최적화)"

        progress_messages.append(await context.bot.send_message(chat_id, "✅ 음성 처리 완료! [100%]"))

        # Send result
        if result:
            final_text = f"🎤 {mode} 처리 결과 ({duration:.1f}초):\n\n{result}"
            await context.bot.send_message(chat_id, final_text)

            # Save to memory
            await save_memory(user_id, username, f"[음성] {duration:.1f}초", result)

    except Exception as e:
        logger.error(f"Voice processing error: {e}")
        error_msg = f"음성 처리 중 오류가 발생했어요: {str(e)[:100]}"
        await context.bot.send_message(chat_id, error_msg)
    finally:
        # Clean up
        try:
            for path in [ogg_path, wav_path]:
                if os.path.exists(path):
                    os.remove(path)
        except Exception:
            pass


async def process_with_gemini_multimodal(ogg_path: str, duration: float, chat_id: int, context, progress_messages):
    """Process short audio with Gemini 2.5 Flash multimodal"""
    # Send progress update
    progress_messages.append(await context.bot.send_message(chat_id, f"🎤 {duration:.1f}초 (짧음) - Gemini 2.5 Flash 멀티모달 분석 중... [60%]"))

    # Upload audio directly to Gemini
    import google.generativeai as genai
    audio_data = open(ogg_path, "rb").read()
    audio_part = {"mime_type": "audio/ogg", "data": audio_data}

    prompt = (
        "이 음성 메시지를 한국어로 전사하고 적절히 요약/답변해주세요.\n"
        "음성 내용에 직접 답할 수 있는 질문이면 답변도 제공해주세요.\n"
        "항상 한국어로만 답변하고, Markdown 표/코드블록 없이 간결한 문장으로 답하세요."
    )

    # Call Gemini in thread pool
    def _call_gemini():
        response = gemini_model.generate_content([prompt, audio_part])
        return response.text.strip()

    result = await asyncio.to_thread(_call_gemini)
    return format_plain(result)


async def process_with_whisper_gemini(ogg_path: str, wav_path: str, duration: float, chat_id: int, context, progress_messages):
    """Process long audio with Whisper + Gemini"""
    # Send progress update
    progress_messages.append(await context.bot.send_message(chat_id, f"🎤 {duration:.1f}초 (김음) - Whisper로 전사 중... [60%]"))

    # Convert ogg to wav (async)
    try:
        proc = await asyncio.create_subprocess_exec(
            "ffmpeg", "-y", "-i", ogg_path, "-ar", "16000", "-ac", "1", wav_path,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        _stdout, _stderr = await proc.communicate()
        if proc.returncode != 0:
            raise RuntimeError("ffmpeg 변환 실패")
    except Exception as e:
        raise Exception(f"오디오 변환 실패: {str(e)}")

    # Send progress update
    progress_messages.append(await context.bot.send_message(chat_id, f"🎤 전사 완료! Gemini로 요약 중... [80%]"))

    # Whisper transcription (in thread pool)
    try:
        from faster_whisper import WhisperModel
        if not hasattr(process_with_whisper_gemini, "_whisper"):
            process_with_whisper_gemini._whisper = WhisperModel("base", device="cpu", compute_type="int8")
        wmodel = process_with_whisper_gemini._whisper

        def _transcribe():
            segs, _info = wmodel.transcribe(wav_path, language="ko", vad_filter=True)
            return " ".join([s.text.strip() for s in segs if s.text]).strip()

        transcription = await asyncio.to_thread(_transcribe)

        if not transcription:
            return "음성에서 텍스트를 인식하지 못했어요."

        # Gemini summary (in thread pool)
        def _summarize():
            prompt = (
                f"다음 음성 메시지가 전사된 텍스트입니다. 적절히 요약하거나 답변해 주세요:\n\n{transcription}\n\n"
                "항상 한국어로만 답변하고, Markdown 표/코드블록 없이 간결한 문장으로 답하세요."
            )
            response = gemini_model.generate_content(prompt)
            return response.text.strip()

        result = await asyncio.to_thread(_summarize)
        return format_plain(result)

    except ImportError:
        return "faster-whisper가 설치되어 있지 않아요. 백엔드 관리자에게 문의해주세요."


async def handle_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    docs = recent_documents.get(user_id, [])[-5:]
    if not docs:
        await reply_text(update, "저장된 최근 문서가 없어요.")
        return
    lines = [f"{i+1}. {d['file_name']} ({d['text_length']}자)" for i, d in enumerate(docs)]
    await reply_text(update, "최근 문서 목록:\n" + "\n".join(lines))


async def handle_mode(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /mode command - Switch user interaction mode"""
    args = context.args
    
    if not args:
        current_mode = user_modes.get(str(update.effective_user.id), 'hybrid')
        mode_names = {
            'beginner': '🔰 초보자 (자연어 우선)',
            'advanced': '⚡ 파워유저 (명령어 우선)',
            'hybrid': '🎭 하이브리드 (자연어 + 명령어)'
        }
        mode_name = mode_names.get(current_mode, 'unknown')
        
        await reply_text(update,
            f"📋 **현재 모드**: {mode_name}\n\n"
            "🎯 **모드 변경**:\n"
            "/mode beginner → 초보자 모드 (자연어 + 상세 설명)\n"
            "/mode advanced → 파워유저 모드 (명령어 + 간결한 응답)\n"
            "/mode hybrid → 하이브리드 모드 (자연어 파악 + 명령어 제안) [기본]\n\n"
            "💡 **예시**:\n"
            "초보자: '메일 확인해줘' → 상세한 응답\n"
            "파워유저: '/gmail_list' → 간결한 응답\n"
            "하이브리드: '메일 확인' → 실행 + 명령어 제안"
        )
        return
    
    mode = args[0].lower()
    if mode not in ['beginner', 'advanced', 'hybrid']:
        await reply_text(update, "❌ 올바른 모드를 선택해주세요: beginner, advanced, hybrid")
        return
    
    user_id = str(update.effective_user.id)
    user_modes[user_id] = mode
    save_user_modes()
    
    mode_names = {
        'beginner': '🔰 초보자 모드',
        'advanced': '⚡ 파워유저 모드',
        'hybrid': '🎭 하이브리드 모드'
    }
    
    await reply_text(update,
        f"✅ **모드가 변경되었습니다!**\n"
        f"→ {mode_names[mode]}\n\n"
        f"이제 '{args[0]}' 스타일로 대화할 수 있어요!"
    )


# ========== Google Drive Sync Handlers ==========

async def handle_drive(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /drive command - show Google Drive sync help"""
    help_text = (
        "📁 **Google Drive 동기화 가이드**\n\n"
        "**사용 가능한 명령어:**\n"
        "• `/drive` - 이 도움말 보기\n"
        "• `/drivelist` - 드라이브 파일 목록 보기\n"
        "• `/driveget <file_id>` - 드라이브에서 파일 가져오기\n"
        "• `/drivesync` - 새로 올라온 파일 확인\n\n"
        "**자동 동기화:**\n"
        "✓ 텔레그램 파일 자동 드라이브 저장 + Gemini 분석\n\n"
        "**지원 파일 형식:**\n"
        "✓ 텍스트: txt, md, py, js, html, css, json, xml, csv 등\n"
        "✓ Office: pdf, docx, pptx, xlsx\n"
        "✓ 압축: zip (내용 미리보기)\n\n"
        "**예시:**\n"
        "1. `/drivelist` - 전체 파일 목록 보기\n"
        "2. `/driveget 1A2B3C4D` - ID가 1A2B3C4D인 파일 다운로드\n"
        "3. `/drivesync` - 새 파일 체크\n"
        "4. 파일 전송 → 자동 드라이브 저장 + 분석\n"
    )
    await reply_text(update, help_text)


async def handle_drive_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /drivelist command - list all files in Google Drive"""
    progress_messages = []
    progress_messages.append(await update.message.reply_text("📁 드라이브 파일 목록 조회 중... [0%]"))

    try:
        # Add backend to path for Telegram handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from services.drive_sync import get_folder_files, format_file_list

        progress_messages.append(await update.message.reply_text("📂 드라이브 연결 중... [30%]"))

        files = get_folder_files()

        progress_messages.append(await update.message.reply_text("📋 파일 목록 생성 중... [70%]"))

        result = format_file_list(files)

        progress_messages.append(await update.message.reply_text("✅ 조회 완료! [100%]"))

        await reply_text(update, result)

    except Exception as e:
        logger.error(f"Drive list error: {e}")
        await reply_text(update, f"드라이브 목록 조회 중 오류가 발생했어요: {str(e)[:100]}")


async def handle_drive_get(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /driveget command - download a file from Google Drive"""
    args = context.args
    if not args:
        await reply_text(update, "사용법: `/driveget <file_id>`\n\n예: `/driveget 1A2B3C4D`")
        return

    file_id = args[0]

    progress_messages = []
    progress_messages.append(await update.message.reply_text(f"📥 드라이브에서 파일 다운로드 중... [0%]"))

    try:
        # Add backend to path for Telegram handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.drive_sync import get_file_info, download_file

        progress_messages.append(await update.message.reply_text("📂 파일 정보 조회 중... [30%]"))

        file_info = get_file_info(file_id)

        if not file_info:
            progress_messages.append(await update.message.reply_text("❌ 파일을 찾을 수 없습니다 [100%]"))
            await reply_text(update, "❌ 파일을 찾을 수 없어요. File ID를 확인해주세요.")
            return

        file_name = file_info['name']
        progress_messages.append(await update.message.reply_text(f"📄 {file_name} 다운로드 중... [60%]"))

        # Download file
        tmp_path = os.path.join(tempfile.gettempdir(), f"drive_download_{file_id}_{file_name}")
        success = download_file(file_id, tmp_path)

        if not success:
            progress_messages.append(await update.message.reply_text("❌ 다운로드 실패 [100%]"))
            await reply_text(update, "❌ 파일 다운로드에 실패했어요.")
            return

        progress_messages.append(await update.message.reply_text("✅ 다운로드 완료! [100%]"))

        # Send file to Telegram
        with open(tmp_path, 'rb') as f:
            from telegram import InputFile
            await context.bot.send_document(
                chat_id=update.effective_chat.id,
                document=InputFile(f, filename=file_name),
                caption=f"📄 **드라이브에서 가져온 파일**: {file_name}"
            )

        # Clean up
        try:
            os.remove(tmp_path)
        except Exception:
            pass

    except Exception as e:
        logger.error(f"Drive get error: {e}")
        await reply_text(update, f"파일 다운로드 중 오류가 발생했어요: {str(e)[:100]}")


async def handle_drive_sync(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /drivesync command - check for new files in Google Drive"""
    progress_messages = []
    progress_messages.append(await update.message.reply_text("🔍 드라이브 새 파일 확인 중... [0%]"))

    try:
        # Add backend to path for Telegram handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.drive_sync import check_new_files, get_folder_files, check_deleted_files

        progress_messages.append(await update.message.reply_text("📂 드라이브 스캔 중... [50%]"))

        # Get current files and check for new/deleted
        current_files = get_folder_files()
        new_files = check_new_files()
        deleted_files = check_deleted_files(current_files)

        progress_messages.append(await update.message.reply_text("✅ 확인 완료! [100%]"))

        # Format results
        result_lines = []
        has_changes = False

        if new_files:
            has_changes = True
            result_lines.append(f"🆕 **새로 올라온 파일** ({len(new_files)}개):\n")
            for i, file in enumerate(new_files, 1):
                file_type = "📁 폴더" if file.get('mimeType') == 'application/vnd.google-apps.folder' else "📄 파일"
                result_lines.append(f"{i}. {file_type}: **{file['name']}**")
                result_lines.append(f"   ID: `{file['id']}`")
            result_lines.append("")

        if deleted_files:
            has_changes = True
            result_lines.append(f"🗑️ **삭제된 파일** ({len(deleted_files)}개):\n")
            for i, file in enumerate(deleted_files, 1):
                result_lines.append(f"{i}. **{file['name']}**")
                result_lines.append(f"   ID: `{file['id']}`")
            result_lines.append("")

        if not has_changes:
            await reply_text(update, "📭 새 파일이 없습니다.")
        else:
            await reply_text(update, "\n".join(result_lines).strip())

    except Exception as e:
        logger.error(f"Drive sync error: {e}")
        await reply_text(update, f"새 파일 확인 중 오류가 발생했어요: {str(e)[:100]}")


# ========== Gmail Handlers ==========

async def handle_gmail_on(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /gmail_on command - Start Gmail monitoring"""
    global gmail_monitoring_state

    if gmail_monitoring_state["enabled"]:
        await reply_text(update,
            "🟡 **Gmail 감시가 이미 실행 중이에요!**\n"
            f"- 현재까지 {gmail_monitoring_state['total_emails']}개 메일 처리됨\n"
            "- `/gmail_status`로 상세 상태 확인")
        return

    # Test Gmail connection
    test_msg = await reply_text(update, "📧 Gmail 연결 테스트 중...")

    try:
        # Add backend to path for Gmail handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.gmail import GmailService

        gmail_service = GmailService()
        if not gmail_service.authenticate():
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=test_msg.message_id,
                text="❌ Gmail 인증 실패. gmail_credentials.json 파일을 확인해주세요.\n\n"
                     "📋 설정 방법:\n"
                     "1. https://console.cloud.google.com/ 접속\n"
                     "2. Gmail API 활성화\n"
                     "3. OAuth 2.0 클라이언트 ID 생성\n"
                     "4. 다운로드한 파일을 gmail_credentials.json으로 저장"
            )
            return

        # Test with 1 email
        test_emails = gmail_service.get_recent_emails(max_results=1)

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=test_msg.message_id,
            text="✅ Gmail 연결 성공! 감시를 시작합니다..."
        )

        # Start monitoring
        gmail_monitoring_state["enabled"] = True
        gmail_monitoring_state["total_emails"] = 0
        gmail_monitoring_state["start_time"] = datetime.now().isoformat()
        start_gmail_monitoring()

        await asyncio.sleep(1)

        final_msg = """
🟢 **Gmail 실시간 감시 시작!**

📋 **감시 설정**:
- 확인 주기: 5분마다
- 대상: 읽지 않은 메일만
- AI 요약: Gemini 2.5 Flash
- 즉시 텔레그램 알림

💡 **명령어**:
- `/gmail_off` - 감시 중지
- `/gmail_status` - 상태 확인
- `/gmail_list` - 최근 메일 목록
        """.strip()

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=test_msg.message_id,
            text=final_msg
        )

    except Exception as e:
        logger.error(f"Gmail start error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=test_msg.message_id,
            text=f"❌ Gmail 연결 실패: {str(e)[:100]}"
        )


async def handle_gmail_off(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /gmail_off command - Stop Gmail monitoring"""
    global gmail_monitoring_state

    if not gmail_monitoring_state["enabled"]:
        await reply_text(update, "🔴 Gmail 감시가 이미 중지되어 있어요!")
        return

    gmail_monitoring_state["enabled"] = False
    total_processed = gmail_monitoring_state.get("total_emails", 0)

    stop_message = f"""
📪 **Gmail 감시 중지됨**

📊 **이번 세션 통계**:
- 처리된 메일: {total_processed}개
- 감시 시간: {gmail_monitoring_state.get('start_time', '확인 불가')}부터

💡 **재시작하려면**:
- `/gmail_on` - 감시 다시 시작
- `/gmail_list` - 수동으로 메일 목록 확인
    """.strip()

    await reply_text(update, stop_message)


async def handle_gmail_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /gmail_status command - Check Gmail monitoring status"""
    global gmail_monitoring_state

    status_icon = "🟢" if gmail_monitoring_state["enabled"] else "🔴"
    status_text = "실행 중" if gmail_monitoring_state["enabled"] else "중지됨"

    last_check = gmail_monitoring_state.get("last_check", "없음")
    total_emails = gmail_monitoring_state.get("total_emails", 0)

    # Get unread count if running
    if gmail_monitoring_state["enabled"]:
        try:
            import sys
            import os
            backend_path = os.path.join(os.path.dirname(__file__))
            if backend_path not in sys.path:
                sys.path.insert(0, backend_path)

            from backend.services.gmail import GmailService
            gmail_service = GmailService()
            if gmail_service.authenticate():
                unread_count = gmail_service.get_unread_count()
            else:
                unread_count = "확인 불가"
        except:
            unread_count = "확인 불가"
    else:
        unread_count = "감시 중지됨"

    status_message = f"""
📊 **Gmail 감시 상태**

{status_icon} **상태**: {status_text}
🕒 **마지막 확인**: {last_check}
📧 **처리된 메일**: {total_emails}개
🔵 **현재 받은편지함**: {unread_count}개

⚙️ **설정**:
- 확인 주기: 5분마다
- 대상: 읽지 않은 메일만
- AI 요약: Gemini 2.5 Flash

💡 **사용 가능한 명령어**:
- `/gmail_on` - 감시 시작
- `/gmail_off` - 감시 중지
- `/gmail_list` - 최근 메일 목록
    """.strip()

    await reply_text(update, status_message)


async def handle_gmail_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /gmail_list command - List recent emails"""
    ack_msg = await reply_text(update, "📧 최근 메일 목록 가져오는 중...")

    try:
        # Add backend to path for Gmail handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.gmail import GmailService

        gmail_service = GmailService()
        if not gmail_service.authenticate():
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="❌ Gmail 인증 실패. 먼저 `/gmail_on`으로 인증해주세요."
            )
            return

        # Get recent emails
        recent_emails = gmail_service.get_recent_emails(max_results=20)

        if not recent_emails:
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="📪 읽지 않은 메일이 없어요."
            )
            return

        # Collect email info
        email_list = []
        for i, email_info in enumerate(recent_emails[:10]):  # Max 10
            email_content = gmail_service.get_email_content(email_info['id'])
            if email_content:
                is_unread = "🔵" if 'UNREAD' in email_info.get('labelIds', []) else "⚪"
                email_list.append(f"""
{i+1}. {is_unread} **{email_content['subject'][:40]}**
   👤 {email_content['sender'][:30]}
   🕒 {email_content['date'][:16]}
                """.strip())

        final_message = f"""
📋 **최근 Gmail 목록** (최대 10개)

{chr(10).join(email_list)}

📊 **요약**:
- 전체 확인된 메일: {len(recent_emails)}개
- 🔵 읽지 않은 메일  ⚪ 읽은 메일
- `/gmail_on`으로 실시간 감시 시작 가능
        """.strip()

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=final_message
        )

    except Exception as e:
        logger.error(f"Gmail list error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ Gmail 목록 조회 중 오류가 발생했어요: {str(e)[:100]}"
        )


# ========== Gmail Reply Handlers ==========

async def handle_gmail_reply(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /gmail_reply command - Generate AI-powered reply"""
    args = context.args

    if not args:
        await reply_text(update,
            "📧 **Gmail 답장 생성**\n\n"
            "사용법:\n"
            "• `/gmail_reply <메일ID>` - 특정 메일에 답장\n"
            "• `/gmail_recent` - 최근 메일 목록에서 선택\n\n"
            "예: `/gmail_reply 1a2b3c4d5e6f`\n"
            "   (메일ID는 `/gmail_list`에서 확인 가능)"
        )
        return

    email_id = args[0]
    ack_msg = await reply_text(update, "📧 메일을 분석하고 답장 초안을 생성 중...")

    try:
        from backend.services.gmail_reply import GmailReplyGenerator

        reply_generator = GmailReplyGenerator()
        if not reply_generator.authenticate():
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="❌ Gmail 인증 실패. gmail_credentials.json 파일을 확인해주세요."
            )
            return

        # Get email content
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text="📧 메일 내용 조회 중... [30%]"
        )

        email_content = reply_generator.get_email_content(email_id)
        if not email_content:
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="❌ 메일을 찾을 수 없습니다. 메일 ID를 확인해주세요."
            )
            return

        # Generate reply draft
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text="✏️ AI가 답장 초안을 작성 중... [70%]"
        )

        reply_draft = reply_generator.generate_reply_draft(email_content, tone="professional")

        if reply_draft:
            # Store reply data for callback
            reply_data = {
                'draft': reply_draft['draft'],
                'tone': reply_draft['tone'],
                'original_subject': reply_draft['original_subject'],
                'original_sender': reply_draft['original_sender'],
                'thread_id': reply_draft['thread_id'],
                'email_id': email_id
            }

            # Save to conversation context
            user_id = str(update.effective_user.id)
            conversation_contexts[user_id] = {
                'pending_reply': reply_data,
                'timestamp': datetime.now().isoformat()
            }
            save_conversation_contexts()

            # Show reply draft
            final_message = f"""
📧 **답장 초안 완성!**

**받은 메일:**
👤 {email_content['sender'][:50]}
📋 {email_content['subject'][:60]}

**답장 초안:**
✏️ {reply_draft['draft']}

어떻게 하시겠어요?
            """.strip()

            from telegram import InlineKeyboardButton, InlineKeyboardMarkup

            keyboard = [
                [
                    InlineKeyboardButton("📤 바로 보내기", callback_data=f"send_reply_{email_id}"),
                    InlineKeyboardButton("✏️ 수정하기", callback_data=f"edit_reply_{email_id}")
                ],
                [
                    InlineKeyboardButton("🔄 다른 톤으로", callback_data=f"regenerate_reply_{email_id}"),
                    InlineKeyboardButton("❌ 취소", callback_data="cancel_reply")
                ]
            ]
            reply_markup = InlineKeyboardMarkup(keyboard)

            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text=final_message,
                reply_markup=reply_markup
            )

        else:
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="❌ 답장 생성에 실패했습니다. 다시 시도해주세요."
            )

    except Exception as e:
        logger.error(f"Gmail reply error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ 오류가 발생했습니다: {str(e)[:100]}"
        )


async def handle_gmail_recent(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /gmail_recent command - Show recent emails for quick reply"""
    ack_msg = await reply_text(update, "📧 최근 메일 목록 가져오는 중...")

    try:
        from backend.services.gmail_reply import GmailReplyGenerator

        reply_generator = GmailReplyGenerator()
        if not reply_generator.authenticate():
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="❌ Gmail 인증 실패. gmail_credentials.json 파일을 확인해주세요."
            )
            return

        recent_emails = reply_generator.find_recent_emails(max_results=10)

        if not recent_emails:
            await context.bot.edit_message_text(
                chat_id=update.effective_chat.id,
                message_id=ack_msg.message_id,
                text="📪 읽지 않은 메일이 없어요."
            )
            return

        # Get email contents
        email_list = []
        for i, email_info in enumerate(recent_emails[:5]):  # Max 5
            email_content = reply_generator.get_email_content(email_info['id'])
            if email_content:
                sender_name = email_content['sender'].split('<')[0].strip() if '<' in email_content['sender'] else email_content['sender']
                email_list.append(
                    f"{i+1}. 📧 **{email_content['subject'][:40]}**\n"
                    f"   👤 {sender_name[:30]} | ID: `{email_info['id'][:12]}`"
                )

        final_message = f"""
📋 **최근 메일 목록** (최대 5개)

{chr(10).join(email_list)}

💡 **답장 방법**: `/gmail_reply <메일ID>`
📝 예시: `/gmail_reply {recent_emails[0]['id'][:12]}`
        """.strip()

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=final_message
        )

    except Exception as e:
        logger.error(f"Gmail recent error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ 오류가 발생했습니다: {str(e)[:100]}"
        )


async def handle_reply_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle reply-related inline keyboard callbacks"""
    query = update.callback_query
    data = query.data
    user_id = str(query.from_user.id)

    await query.answer()

    try:
        if data == "cancel_reply":
            if user_id in conversation_contexts:
                del conversation_contexts[user_id]
                save_conversation_contexts()
            await query.edit_message_text("❌ 답장 작성이 취소되었습니다.")
            return

        # Extract email ID from callback data
        email_id = data.split('_')[-1]

        # Get pending reply from context
        if user_id not in conversation_contexts or 'pending_reply' not in conversation_contexts[user_id]:
            await query.edit_message_text("❌ 답장 정보를 찾을 수 없습니다. 다시 시도해주세요.")
            return

        reply_data = conversation_contexts[user_id]['pending_reply']

        if data.startswith("send_reply_"):
            await query.edit_message_text("📤 답장을 보내는 중...")

            from backend.services.gmail_reply import GmailReplyGenerator
            reply_generator = GmailReplyGenerator()
            reply_generator.authenticate()

            result = reply_generator.send_reply_email(reply_data)

            if result:
                await query.edit_message_text("✅ 답장을 성공적으로 보냈습니다!")
                # Mark email as read
                reply_generator.mark_as_read(email_id)
            else:
                await query.edit_message_text("❌ 답장 전송에 실패했습니다.")

            # Clear context
            del conversation_contexts[user_id]
            save_conversation_contexts()

        elif data.startswith("edit_reply_"):
            await query.edit_message_text(
                "✏️ 수정하실 내용을 다음 메시지로 보내주세요.\n"
                "기존 답장:\n"
                f"{reply_data['draft']}"
            )

            # Store editing state
            conversation_contexts[user_id]['editing_reply'] = True
            save_conversation_contexts()

        elif data.startswith("regenerate_reply_"):
            await query.edit_message_text("🔄 다른 톤으로 답장을 다시 작성 중...")

            # Regenerate with different tone
            from backend.services.gmail_reply import GmailReplyGenerator
            reply_generator = GmailReplyGenerator()
            reply_generator.authenticate()

            # Get email content again
            email_content = reply_generator.get_email_content(email_id)
            if email_content:
                # Try different tone
                tones = ["friendly", "concise", "professional"]
                current_tone = reply_data.get('tone', 'professional')
                next_tone = tones[(tones.index(current_tone) + 1) % len(tones)]

                new_draft = reply_generator.generate_reply_draft(email_content, tone=next_tone)
                if new_draft:
                    # Update reply data
                    reply_data['draft'] = new_draft['draft']
                    reply_data['tone'] = new_draft['tone']
                    conversation_contexts[user_id]['pending_reply'] = reply_data
                    save_conversation_contexts()

                    # Show new draft
                    final_message = f"""
📧 **새 답장 초안** ({next_tone} 톤)

✏️ {new_draft['draft']}

어떻게 하시겠어요?
                    """.strip()

                    from telegram import InlineKeyboardButton, InlineKeyboardMarkup

                    keyboard = [
                        [
                            InlineKeyboardButton("📤 바로 보내기", callback_data=f"send_reply_{email_id}"),
                            InlineKeyboardButton("✏️ 수정하기", callback_data=f"edit_reply_{email_id}")
                        ],
                        [
                            InlineKeyboardButton("🔄 다른 톤으로", callback_data=f"regenerate_reply_{email_id}"),
                            InlineKeyboardButton("❌ 취소", callback_data="cancel_reply")
                        ]
                    ]
                    reply_markup = InlineKeyboardMarkup(keyboard)

                    await query.edit_message_text(
                        text=final_message,
                        reply_markup=reply_markup
                    )

    except Exception as e:
        logger.error(f"Reply callback error: {e}")
        await query.edit_message_text(f"❌ 오류가 발생했습니다: {str(e)[:100]}")


# ========== Gmail Monitoring Functions ==========

def start_gmail_monitoring():
    """Start Gmail monitoring in background thread"""
    import threading
    if gmail_monitoring_state["thread"] and gmail_monitoring_state["thread"].is_alive():
        return

    gmail_monitoring_state["thread"] = threading.Thread(
        target=gmail_monitor_loop,
        daemon=True
    )
    gmail_monitoring_state["thread"].start()
    logger.info("📧 Gmail monitoring started")


def gmail_monitor_loop():
    """Background Gmail monitoring loop"""
    import time

    try:
        from backend.services.gmail import GmailService
        gmail_service = GmailService()

        if not gmail_service.authenticate():
            logger.error("Gmail authentication failed")
            return

        logger.info("📧 Gmail monitoring worker started")
        gmail_monitoring_state["enabled"] = True
        gmail_monitoring_state["auth_mode"] = getattr(gmail_service, "auth_mode", "unknown")
        gmail_monitoring_state["interval"] = GMAIL_MONITOR_INTERVAL
        if not gmail_monitoring_state.get("start_time"):
            gmail_monitoring_state["start_time"] = datetime.now().isoformat()

        while gmail_monitoring_state.get("enabled", True):
            try:
                logger.info("📧 Checking for new emails...")

                recent_emails = gmail_service.get_recent_emails(max_results=20)
                new_emails = []

                for email_info in recent_emails:
                    email_id = email_info['id']

                    if email_id not in gmail_service.processed_emails:
                        email_content = gmail_service.get_email_content(email_id)
                        if email_content:
                            new_emails.append(email_content)
                            gmail_service.processed_emails.add(email_id)
                            try:
                                gmail_service.mark_as_read(email_id)
                            except Exception as mark_err:
                                logger.warning(f"Failed to mark email as read ({email_id}): {mark_err}")

                if new_emails:
                    logger.info(f"📧 Found {len(new_emails)} new emails")
                    gmail_monitoring_state["total_emails"] += len(new_emails)

                    for email_data in new_emails:
                        gmail_monitoring_state["recent_emails"].appendleft({
                            "subject": email_data.get('subject', '')[:120],
                            "sender": email_data.get('sender', '')[:120],
                            "date": email_data.get('date', ''),
                            "processed_at": datetime.utcnow().isoformat() + 'Z'
                        })
                        usage_metrics["gmail_notifications"] += 1
                        asyncio.run_coroutine_threadsafe(
                            process_and_send_email(email_data),
                            asyncio.get_event_loop()
                        )

                gmail_service.save_processed_emails()
                gmail_monitoring_state["last_check"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

                interval = max(30, gmail_monitoring_state.get("interval", GMAIL_MONITOR_INTERVAL))
                for _ in range(interval):
                    if not gmail_monitoring_state.get("enabled", True):
                        break
                    time.sleep(1)

            except Exception as e:
                logger.error(f"Gmail monitoring error: {e}")
                time.sleep(60)

        logger.info("📧 Gmail monitoring worker stopped")

    except Exception as e:
        logger.error(f"Gmail monitoring loop error: {e}")
    finally:
        gmail_monitoring_state["enabled"] = False


async def process_and_send_email(email_data):
    """Process email with Gemini and send to Telegram"""
    if not gemini_model:
        return

    try:
        # Create progress message
        progress_msg = await _app_instance.bot.send_message(
            chat_id=_app_instance.chat_ids[0] if _app_instance.chat_ids else None,
            text="📧 새 메일 분석 중..."
        )

        # Gemini summarization
        prompt = f"""
        다음 이메일을 한국어로 요약해주세요:

        보낸사람: {email_data['sender']}
        제목: {email_data['subject']}
        내용: {email_data['body']}

        요약 형식:
        - 핵심 내용 (2-3문장)
        - 중요도 (높음/보통/낮음)
        - 필요한 액션이 있다면 언급
        """

        response = gemini_model.generate_content(prompt)
        summary = format_plain(response.text)

        # Final message
        final_message = f"""
📧 **새 메일 요약**

👤 **보낸사람**: {email_data['sender']}
📝 **제목**: {email_data['subject']}
🕒 **시간**: {email_data['date']}

🤖 **AI 요약**:
{summary}
        """.strip()

        await _app_instance.bot.edit_message_text(
            chat_id=progress_msg.chat_id,
            message_id=progress_msg.message_id,
            text=final_message
        )

    except Exception as e:
        logger.error(f"Email processing error: {e}")


# ========== Calendar Handlers ==========

async def handle_cal_on(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_on command - Start Calendar monitoring"""
    global calendar_monitoring_state

    if calendar_monitoring_state["enabled"]:
        await reply_text(update,
            "🟡 **Calendar 감시가 이미 실행 중이에요!**\n"
            f"- 현재까지 {calendar_monitoring_state['total_alerts']}개 알림 보냄\n"
            "- `/cal_status`로 상세 상태 확인")
        return

    # Test Calendar connection
    test_msg = await reply_text(update, "🗓️ Calendar 연결 테스트 중...")

    try:
        # Add backend to path for Calendar handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.calendar import get_calendar_service

        # Test Calendar connection
        calendar_service = get_calendar_service()
        test_events = calendar_service.get_today_events()

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=test_msg.message_id,
            text="✅ Calendar 연결 성공! 감시를 시작합니다..."
        )

        # Start monitoring
        calendar_monitoring_state["enabled"] = True
        calendar_monitoring_state["total_alerts"] = 0
        calendar_monitoring_state["start_time"] = datetime.now().isoformat()
        calendar_monitoring_state["alerted_events"] = set()
        start_calendar_monitoring()

        await asyncio.sleep(1)

        final_msg = """
🟢 **Calendar 실시간 감시 시작!**

📋 **감시 설정**:
- 확인 주기: 5분마다
- 대상: 다가오는 일정 (30분 전 알림)
- AI 분석: Gemini 2.5 Flash
- 즉시 텔레그램 알림

💡 **명령어**:
- `/cal_off` - 감시 중지
- `/cal_status` - 상태 확인
- `/cal_today` - 오늘 일정
- `/cal_tomorrow` - 내일 일정
- `/cal_week` - 이번 주 일정
- `/cal_search <키워드>` - 일정 검색
        """.strip()

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=test_msg.message_id,
            text=final_msg
        )

    except Exception as e:
        logger.error(f"Calendar start error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=test_msg.message_id,
            text=f"❌ Calendar 연결 실패: {str(e)[:100]}"
        )


async def handle_cal_off(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_off command - Stop Calendar monitoring"""
    global calendar_monitoring_state

    if not calendar_monitoring_state["enabled"]:
        await reply_text(update, "🔴 Calendar 감시가 이미 중지되어 있어요!")
        return

    calendar_monitoring_state["enabled"] = False
    total_alerts = calendar_monitoring_state.get("total_alerts", 0)

    stop_message = f"""
📅 **Calendar 감시 중지됨**

📊 **이번 세션 통계**:
- 보낸 알림: {total_alerts}개
- 감시 시간: {calendar_monitoring_state.get('start_time', '확인 불가')}부터

💡 **재시작하려면**:
- `/cal_on` - 감시 다시 시작
- `/cal_today` - 수동으로 오늘 일정 확인
    """.strip()

    await reply_text(update, stop_message)


async def handle_cal_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_status command - Check Calendar monitoring status"""
    global calendar_monitoring_state

    status_icon = "🟢" if calendar_monitoring_state["enabled"] else "🔴"
    status_text = "실행 중" if calendar_monitoring_state["enabled"] else "중지됨"

    last_check = calendar_monitoring_state.get("last_check", "없음")
    total_alerts = calendar_monitoring_state.get("total_alerts", 0)

    # Get today's events if running
    if calendar_monitoring_state["enabled"]:
        try:
            import sys
            import os
            backend_path = os.path.join(os.path.dirname(__file__))
            if backend_path not in sys.path:
                sys.path.insert(0, backend_path)

            from backend.services.calendar import get_calendar_service
            calendar_service = get_calendar_service()
            today_events = calendar_service.get_today_events()
            today_count = len(today_events)
        except:
            today_count = "확인 불가"
    else:
        today_count = "감시 중지됨"

    status_message = f"""
📊 **Calendar 감시 상태**

{status_icon} **상태**: {status_text}
🕒 **마지막 확인**: {last_check}
📅 **보낸 알림**: {total_alerts}개
📋 **오늘 일정**: {today_count}개

⚙️ **설정**:
- 확인 주기: 5분마다
- 알림: 30분 전 일정
- AI 분석: Gemini 2.5 Flash

💡 **사용 가능한 명령어**:
- `/cal_on` - 감시 시작
- `/cal_off` - 감시 중지
- `/cal_today` - 오늘 일정
- `/cal_tomorrow` - 내일 일정
- `/cal_week` - 이번 주 일정
- `/cal_search <키워드>` - 일정 검색
    """.strip()

    await reply_text(update, status_message)


async def handle_cal_today(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_today command - Show today's events"""
    ack_msg = await reply_text(update, "🗓️ 오늘 일정 조회 중...")

    try:
        # Add backend to path for Calendar handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.calendar import get_calendar_service, format_event_list

        calendar_service = get_calendar_service()
        today_events = calendar_service.get_today_events()

        result = format_event_list(today_events, "오늘의 일정")

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=result
        )

    except Exception as e:
        logger.error(f"Calendar today error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ 오늘 일정 조회 중 오류가 발생했어요: {str(e)[:100]}"
        )


async def handle_cal_tomorrow(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_tomorrow command - Show tomorrow's events"""
    ack_msg = await reply_text(update, "🗓️ 내일 일정 조회 중...")

    try:
        # Add backend to path for Calendar handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.calendar import get_calendar_service, format_event_list

        calendar_service = get_calendar_service()
        tomorrow_events = calendar_service.get_tomorrow_events()

        result = format_event_list(tomorrow_events, "내일의 일정")

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=result
        )

    except Exception as e:
        logger.error(f"Calendar tomorrow error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ 내일 일정 조회 중 오류가 발생했어요: {str(e)[:100]}"
        )


async def handle_cal_week(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_week command - Show this week's events"""
    ack_msg = await reply_text(update, "🗓️ 이번 주 일정 조회 중...")

    try:
        # Add backend to path for Calendar handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.calendar import get_calendar_service, format_event_list

        calendar_service = get_calendar_service()
        week_events = calendar_service.get_week_events()

        result = format_event_list(week_events, "이번 주 일정")

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=result
        )

    except Exception as e:
        logger.error(f"Calendar week error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ 이번 주 일정 조회 중 오류가 발생했어요: {str(e)[:100]}"
        )


async def handle_cal_search(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_search command - Search for events"""
    args = context.args
    if not args:
        await reply_text(update, "사용법: `/cal_search <검색어>`\n\n예: `/cal_search 미팅`")
        return

    search_query = " ".join(args)
    ack_msg = await reply_text(update, f"🔍 '{search_query}' 일정 검색 중...")

    try:
        # Add backend to path for Calendar handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.calendar import get_calendar_service, format_event_list

        calendar_service = get_calendar_service()
        search_results = calendar_service.search_events(search_query, max_results=20)

        result = format_event_list(search_results, f"검색 결과: {search_query}")

        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=result
        )

    except Exception as e:
        logger.error(f"Calendar search error: {e}")
        await context.bot.edit_message_text(
            chat_id=update.effective_chat.id,
            message_id=ack_msg.message_id,
            text=f"❌ 일정 검색 중 오류가 발생했어요: {str(e)[:100]}"
        )


# ========== Calendar Monitoring Functions ==========

def start_calendar_monitoring():
    """Start Calendar monitoring in background thread"""
    import threading
    if calendar_monitoring_state["thread"] and calendar_monitoring_state["thread"].is_alive():
        return

    calendar_monitoring_state["thread"] = threading.Thread(
        target=calendar_monitor_loop,
        daemon=True
    )
    calendar_monitoring_state["thread"].start()
    logger.info("🗓️ Calendar monitoring started")


def calendar_monitor_loop():
    """Background Calendar monitoring loop"""
    import time

    try:
        # Add backend to path for Thread
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.calendar import get_calendar_service, get_upcoming_events

        calendar_service = get_calendar_service()

        logger.info("🗓️ Calendar monitoring worker started")

        while calendar_monitoring_state["enabled"]:
            try:
                logger.info("🗓️ Checking for upcoming events...")

                # Get events in next 30 minutes
                upcoming_events = get_upcoming_events(minutes_ahead=30)
                new_alerts = []

                for event in upcoming_events:
                    event_id = event.get('id', '')

                    # Check if already alerted
                    if event_id and event_id not in calendar_monitoring_state["alerted_events"]:
                        new_alerts.append(event)
                        calendar_monitoring_state["alerted_events"].add(event_id)

                # Send notifications for new alerts
                if new_alerts:
                    logger.info(f"🗓️ Found {len(new_alerts)} upcoming events")
                    calendar_monitoring_state["total_alerts"] += len(new_alerts)

                    for event_data in new_alerts:
                        asyncio.run_coroutine_threadsafe(
                            process_and_send_calendar_alert(event_data),
                            asyncio.get_event_loop()
                        )

                calendar_monitoring_state["last_check"] = datetime.now().strftime("%H:%M:%S")

                # Wait 5 minutes
                for _ in range(300):  # Check every second for shutdown
                    if not calendar_monitoring_state["enabled"]:
                        break
                    time.sleep(1)

            except Exception as e:
                logger.error(f"Calendar monitoring error: {e}")
                time.sleep(60)  # Wait 1 minute on error

        logger.info("🗓️ Calendar monitoring worker stopped")

    except Exception as e:
        logger.error(f"Calendar monitoring loop error: {e}")


async def process_and_send_calendar_alert(event_data):
    """Process event and send alert to Telegram"""
    try:
        # Get start and end time
        start = event_data.get('start', {})
        end = event_data.get('end', {})
        
        # Format time
        time_str = ""
        if 'dateTime' in start:
            start_dt = datetime.fromisoformat(start['dateTime'].replace('Z', '+00:00'))
            end_dt = datetime.fromisoformat(end['dateTime'].replace('Z', '+00:00'))
            time_str = f"{start_dt.strftime('%H:%M')} - {end_dt.strftime('%H:%M')}"
        else:
            time_str = "종일"

        title = event_data.get('summary', '제목 없음')
        location = event_data.get('location', '')
        description = event_data.get('description', '')

        # Create message
        alert_message = f"""
🔔 **30분 후 일정 알림**

📅 **일정**: {title}
⏰ **시간**: {time_str}
        """.strip()

        if location:
            alert_message += f"\n📍 **장소**: {location}"

        if description:
            desc_preview = description[:100]
            if len(description) > 100:
                desc_preview += "..."
            alert_message += f"\n📝 **설명**: {desc_preview}"

        alert_message += "\n\n⏰ 준비하세요!"

        # Send to all active chats (for now, broadcast to first chat)
        if _app_instance and _app_instance.chat_ids:
            for chat_id in _app_instance.chat_ids:
                try:
                    await _app_instance.bot.send_message(
                        chat_id=chat_id,
                        text=alert_message
                    )
                except Exception as e:
                    logger.error(f"Failed to send calendar alert to {chat_id}: {e}")

    except Exception as e:
        logger.error(f"Calendar alert processing error: {e}")


async def monitor_drive_changes():
    """Background task to monitor Google Drive for changes"""
    logger.info("🔍 Drive monitoring worker started")

    # Add backend to path for Thread
    import sys
    import os
    backend_path = os.path.join(os.path.dirname(__file__))
    if backend_path not in sys.path:
        sys.path.insert(0, backend_path)

    drive_monitoring_state["enabled"] = True
    if not drive_monitoring_state.get("start_time"):
        drive_monitoring_state["start_time"] = datetime.now().isoformat()
    drive_monitoring_state["interval"] = DRIVE_MONITOR_INTERVAL

    while drive_monitoring_state.get("enabled", True):
        try:
            if not ENABLE_DRIVE_MONITORING:
                await asyncio.sleep(60)
                continue

            from services.drive_sync import (
                get_folder_files, check_new_files, check_deleted_files,
                cache_current_files, load_cached_files
            )

            current_files = get_folder_files()
            drive_monitoring_state["total_files"] = len(current_files)
            drive_monitoring_state["last_check"] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

            deleted_files = check_deleted_files(current_files)
            new_files = check_new_files()

            if new_files or deleted_files:
                summary_parts = []
                if new_files:
                    summary_parts.append(f"🆕 새 파일 {len(new_files)}개")
                if deleted_files:
                    summary_parts.append(f"🗑️ 삭제 {len(deleted_files)}개")
                summary = ", ".join(summary_parts) if summary_parts else "변경 없음"

                logger.info(f"Drive changes detected: {len(new_files)} new, {len(deleted_files)} deleted")
                drive_monitoring_state["recent_changes"].appendleft({
                    "timestamp": datetime.utcnow().isoformat() + 'Z',
                    "new": len(new_files),
                    "deleted": len(deleted_files),
                    "summary": summary
                })
                usage_metrics["drive_sync_runs"] += 1

            if not load_cached_files():
                cache_current_files(current_files)
                logger.info("Initialized Drive file cache")

        except Exception as e:
            logger.error(f"Drive monitoring error: {e}")

        interval = max(60, drive_monitoring_state.get("interval", DRIVE_MONITOR_INTERVAL))
        try:
            for _ in range(interval):
                if not drive_monitoring_state.get("enabled", True):
                    break
                await asyncio.sleep(1)
        except asyncio.CancelledError:
            break

    drive_monitoring_state["enabled"] = False
    logger.info("🔍 Drive monitoring worker stopped")


def extract_text_from_file(file_path: str, file_name: str) -> str:
    """
    Extract text from various file formats
    Supports: txt, md, py, js, html, css, json, xml, csv, pdf, docx, pptx, xlsx, zip
    """
    import os
    import zipfile
    import chardet

    file_ext = os.path.splitext(file_name)[1].lower()

    try:
        # 1. Text-based files (most common)
        if file_ext in ['.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.html',
                        '.htm', '.css', '.scss', '.sass', '.less', '.json', '.xml',
                        '.csv', '.tsv', '.yaml', '.yml', '.ini', '.cfg', '.conf',
                        '.log', '.sql', '.sh', '.bat', '.ps1', '.dockerfile',
                        '.gitignore', '.env', '.properties', '.toml', '.r', '.R']:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                enc = chardet.detect(raw_data).get('encoding') or 'utf-8'
                return raw_data.decode(enc, errors='ignore')

        # 2. PDF files
        elif file_ext == '.pdf':
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                return text
            except ImportError:
                return "PDF 읽기를 위한 PyPDF2가 설치되어 있지 않습니다."

        # 3. Word documents (.docx)
        elif file_ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            except ImportError:
                return "DOCX 읽기를 위한 python-docx가 설치되어 있지 않습니다."

        # 4. PowerPoint (.pptx)
        elif file_ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except ImportError:
                return "PPTX 읽기를 위한 python-pptx가 설치되어 있지 않습니다."

        # 5. Excel files (.xlsx, .xls)
        elif file_ext in ['.xlsx', '.xls']:
            try:
                import pandas as pd
                # Try to read all sheets
                df = pd.read_excel(file_path, sheet_name=None)
                text = ""
                for sheet_name, sheet_df in df.items():
                    text += f"\n=== Sheet: {sheet_name} ===\n"
                    text += sheet_df.to_string(index=False) + "\n"
                return text
            except ImportError:
                return "Excel 읽기를 위한 pandas가 설치되어 있지 않습니다."

        # 6. ZIP archives (extract and read text files inside)
        elif file_ext == '.zip':
            try:
                text = "=== ZIP 아카이브 내용 ===\n"
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    file_list = zip_ref.namelist()
                    text += f"총 {len(file_list)}개 파일\n\n"
                    for file_in_zip in file_list[:10]:  # Show first 10 files
                        text += f"• {file_in_zip}\n"
                        # If it's a text file, try to extract and read
                        if any(file_in_zip.lower().endswith(ext) for ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml']):
                            try:
                                content = zip_ref.read(file_in_zip).decode('utf-8', errors='ignore')
                                text += f"  내용 미리보기:\n{content[:500]}...\n"
                            except:
                                pass
                    if len(file_list) > 10:
                        text += f"\n... 외 {len(file_list) - 10}개 파일"
                return text
            except Exception as e:
                return f"ZIP 파일 읽기 오류: {str(e)}"

        # 7. Other binary files
        else:
            return f"지원하지 않는 파일 형식: {file_ext}\n파일 크기: {os.path.getsize(file_path)} bytes"

    except Exception as e:
        return f"파일 읽기 오류: {str(e)}"


async def handle_document_auto_save(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Auto-save all documents to Google Drive and analyze with Gemini"""
    doc = update.message.document
    if not doc:
        return

    progress_messages = []
    progress_messages.append(await update.message.reply_text(f"📁 {doc.file_name} Google Drive 자동 저장 중... [0%]"))

    file = await context.bot.get_file(doc.file_id)
    tmp = os.path.join(tempfile.gettempdir(), f"{doc.file_id}_{doc.file_name}")

    doc_indicator = ActionIndicator(context, update.effective_chat.id, ChatAction.UPLOAD_DOCUMENT)
    await doc_indicator.__aenter__()
    await file.download_to_drive(tmp)

    progress_messages.append(await update.message.reply_text("📁 파일 다운로드 완료. 드라이브 저장 중... [30%]"))

    try:
        # Add backend to path for Telegram handlers
        import sys
        import os
        backend_path = os.path.join(os.path.dirname(__file__))
        if backend_path not in sys.path:
            sys.path.insert(0, backend_path)

        from backend.services.drive_sync import upload_file

        # Upload to Google Drive
        result = upload_file(tmp)

        if result:
            progress_messages.append(await update.message.reply_text("✅ Google Drive 저장 완료! [100%]"))

            file_id = result.get('id', 'N/A')
            web_link = result.get('webViewLink', '')

            # Send confirmation
            confirm_text = (
                f"✅ **{doc.file_name}** Google Drive에 자동 저장되었습니다!\n\n"
                f"📋 파일 ID: `{file_id}`"
            )
            if web_link:
                confirm_text += f"\n🔗 [드라이브에서 보기]({web_link})"

            # Add smart suggestions based on file type
            file_ext = os.path.splitext(doc.file_name)[1].lower()
            suggestions = await suggest_smart_actions(file_ext, doc.file_name)
            confirm_text += f"\n\n{suggestions}"

            await reply_text(update, confirm_text)

            # Analyze with Gemini if GEMINI is available
            if GEMINI_API_KEY and gemini_model:
                try:
                    progress_messages.append(await update.message.reply_text("🧠 Gemini 문서 분석 중... [70%]"))

                    # Extract text based on file type
                    extracted_text = extract_text_from_file(tmp, doc.file_name)

                    if extracted_text and len(extracted_text.strip()) > 0:
                        prompt = f"다음 문서를 요약/분석해줘. 파일명: {doc.file_name}\n\n{extracted_text}"
                        prompt += "\n\n항상 한국어로만 답변하고, Markdown 표/코드블록 없이 간결한 문장으로 답하세요."

                        def _call_gemini_doc():
                            resp = gemini_model.generate_content(prompt)
                            return resp.text.strip()

                        answer = await asyncio.to_thread(_call_gemini_doc)
                        answer = format_plain(answer)

                        analysis_text = f"\n\n📄 **문서 분석 결과**:\n\n{answer}"
                        await reply_text(update, analysis_text)
                    else:
                        logger.warning(f"No text extracted from {doc.file_name}")

                except Exception as e:
                    logger.error(f"Document analysis error: {e}")
                    # Don't fail the upload if analysis fails

        else:
            progress_messages.append(await update.message.reply_text("❌ 드라이브 저장 실패 [100%]"))
            await reply_text(update, "❌ Google Drive 저장에 실패했어요. 권한을 확인해주세요.")

    except Exception as e:
        logger.error(f"Auto-save error: {e}")
        await reply_text(update, f"자동 저장 중 오류가 발생했어요: {str(e)[:100]}")
    finally:
        try:
            os.remove(tmp)
        except Exception:
            pass
        await doc_indicator.__aexit__(None, None, None)


def main():
    print("=== 125 Unified Telegram Bot (Gemini 2.5 Flash + Drive Sync) ===")
    print(f"TELEGRAM_BOT_TOKEN: {'Set' if TELEGRAM_BOT_TOKEN else 'Not Found'}")
    print(f"GEMINI_API_KEY: {'Set' if GEMINI_API_KEY else 'Not Found'}")
    print(f"Supabase: {'Set' if (SUPABASE_URL and SUPABASE_KEY) else 'Not Set'}")
    print(f"Google Drive: {'Set' if os.path.exists(os.path.join(os.path.dirname(__file__), '..', 'service_account.json')) else 'Not Set'}")
    print(f"Drive Monitoring: {'Enabled' if ENABLE_DRIVE_MONITORING else 'Disabled'} (interval: {DRIVE_MONITOR_INTERVAL}s)")

    if not TELEGRAM_BOT_TOKEN:
        print("ERROR: TELEGRAM_BOT_TOKEN is missing")
        return

    app = Application.builder().token(TELEGRAM_BOT_TOKEN).build()

    # Store app instance for Drive monitoring
    global _app_instance
    _app_instance = app

    app.add_handler(CommandHandler("start", handle_start))
    app.add_handler(CommandHandler("list", handle_list))
    app.add_handler(CommandHandler("mode", handle_mode))
    app.add_handler(CommandHandler("drive", handle_drive))
    app.add_handler(CommandHandler("drivelist", handle_drive_list))
    app.add_handler(CommandHandler("driveget", handle_drive_get))
    app.add_handler(CommandHandler("drivesync", handle_drive_sync))
    app.add_handler(CommandHandler("gmail_on", handle_gmail_on))
    app.add_handler(CommandHandler("gmail_off", handle_gmail_off))
    app.add_handler(CommandHandler("gmail_status", handle_gmail_status))
    app.add_handler(CommandHandler("gmail_list", handle_gmail_list))
    app.add_handler(CommandHandler("cal_on", handle_cal_on))
    app.add_handler(CommandHandler("cal_off", handle_cal_off))
    app.add_handler(CommandHandler("cal_status", handle_cal_status))
    app.add_handler(CommandHandler("cal_today", handle_cal_today))
    app.add_handler(CommandHandler("cal_tomorrow", handle_cal_tomorrow))
    app.add_handler(CommandHandler("cal_week", handle_cal_week))
    app.add_handler(CommandHandler("cal_search", handle_cal_search))
    app.add_handler(CommandHandler("gmail_reply", handle_gmail_reply))
    app.add_handler(CommandHandler("gmail_recent", handle_gmail_recent))

    # Add callback query handler for inline keyboards
    app.add_handler(MessageHandler(filters.CALLBACK_QUERY, handle_callback_query))

    app.add_handler(MessageHandler(filters.Document.ALL, handle_document_auto_save))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    app.add_handler(MessageHandler(filters.VOICE, handle_voice))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

    logger.info("Handlers registered. Starting polling...")

    # Start Drive monitoring in a separate thread
    if ENABLE_DRIVE_MONITORING:
        def run_monitoring():
            import asyncio
            # Create new event loop for the thread
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            loop.run_until_complete(monitor_drive_changes())

        import threading
        monitor_thread = threading.Thread(target=run_monitoring, daemon=True)
        monitor_thread.start()
        logger.info("Drive monitoring worker started in background thread")

    app.run_polling()


if __name__ == "__main__":
    main()
