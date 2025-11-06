#!/usr/bin/env python3
"""
125 Build Automation - Telegram Bot Runner (Gemini 2.0 Flash Multimodal)
- Single file handling text/document/image/voice with Gemini 2.0 Flash
- Free chat with memory (Supabase optional)
- Document/Image/Voice processed directly with Gemini's multimodal capabilities
- Google Drive bidirectional sync
"""
import os
import sys
import logging
from datetime import datetime
from typing import Dict, List, Any
import tempfile
import asyncio

# Add current directory to Python path
current_dir = os.getcwd()
if current_dir not in sys.path:
    sys.path.insert(0, current_dir)

from dotenv import load_dotenv

# Ensure .env is loaded from the backend directory regardless of package depth.
load_dotenv(os.path.join(os.path.dirname(__file__), "..", "..", ".env"))

TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_ANON_KEY")

# logging
from backend.core import build_application
from backend.utils.logger import configure_logging, LOG_DIR
from backend.bots.main.services.drive import handlers as drive_handlers
from backend.bots.main.services.calendar import handlers as calendar_handlers
from backend.bots.main.services.media import handlers as media_handlers
from backend.bots.main.services.text import handlers as text_handlers

configure_logging()
logger = logging.getLogger("unified_bot")

logger.info(f"Logging to directory: {LOG_DIR.resolve()}")

RUNTIME = sys.modules[__name__]

# Disable httpx logging to prevent token exposure
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("telegram").setLevel(logging.WARNING)

# telegram
try:
    from telegram import Update
    from telegram.constants import ChatAction
    from telegram.ext import ContextTypes
except ImportError:
    logger.error("python-telegram-bot is not installed. pip install python-telegram-bot==21.6")
    sys.exit(1)

# gemini (multimodal)
gemini_model = None
if GEMINI_API_KEY:
    try:
        import google.generativeai as genai
        genai.configure(api_key=GEMINI_API_KEY)
        gemini_model = genai.GenerativeModel('gemini-2.5-flash')
        logger.info("Using Gemini 2.5 Flash (multimodal)")
    except Exception as e:
        logger.error(f"Gemini setup failed: {e}")
else:
    logger.warning("GEMINI_API_KEY not set; chat will be disabled")

# supabase (optional memory)
supabase = None
if SUPABASE_URL and SUPABASE_KEY:
    try:
        from supabase import create_client
        supabase = create_client(SUPABASE_URL, SUPABASE_KEY)
    except Exception as e:
        logger.warning(f"Supabase init failed: {e}")

# in-memory recent docs (fallback)
recent_documents: Dict[int, List[Dict[str, Any]]] = {}

# Smart audio processing configuration
SHORT_AUDIO_THRESHOLD = int(os.getenv("SHORT_AUDIO_THRESHOLD", "30"))  # 30초 이하
LONG_AUDIO_THRESHOLD = int(os.getenv("LONG_AUDIO_THRESHOLD", "300"))  # 5분 이상
MID_LENGTH_MODEL = os.getenv("MID_LENGTH_AUDIO", "gemini")  # 30초-5분 기본

# Drive monitoring configuration
DRIVE_MONITOR_INTERVAL = int(os.getenv("DRIVE_MONITOR_INTERVAL", "300"))  # 5분 (300초)
ENABLE_DRIVE_MONITORING = os.getenv("ENABLE_DRIVE_MONITORING", "true").lower() == "true"

# Global application instance for Drive monitoring
_app_instance = None

# Drive monitoring state control
drive_monitoring_state = {
    "enabled": False,
    "thread": None,
    "last_check": None,
    "total_files": 0,
    "start_time": None
}

# Gmail monitoring state control
gmail_monitoring_state = {
    "enabled": False,
    "thread": None,
    "last_check": None,
    "total_emails": 0,
    "start_time": None
}

# Calendar monitoring state control
calendar_monitoring_state = {
    "enabled": False,
    "thread": None,
    "last_check": None,
    "total_alerts": 0,
    "start_time": None,
    "alerted_events": set()  # Track alerted event IDs
}


def get_audio_duration(ogg_path: str) -> float:
    """Get audio duration in seconds using ffprobe (if available) or estimate"""
    try:
        import subprocess
        result = subprocess.run(
            ["ffprobe", "-v", "quiet", "-show_entries", "format=duration",
             "-of", "csv=p=0", ogg_path],
            capture_output=True, text=True, check=True
        )
        return float(result.stdout.strip())
    except Exception:
        # Fallback: estimate based on file size (rough)
        # ~1MB per minute at 64kbps
        size_mb = os.path.getsize(ogg_path) / (1024 * 1024)
        return size_mb * 60 * 0.7  # conservative estimate


def format_plain(text: str, max_len: int = 1200) -> str:
    """Format Gemini response to Telegram-friendly plain text"""
    import re
    # Remove code blocks
    text = re.sub(r"```.*?```", "", text, flags=re.DOTALL)
    # Remove tables
    text = re.sub(r"\|.*\|", "", text)
    # Remove header symbols (keep line breaks)
    text = re.sub(r"^#{1,6}\s*", "", text, flags=re.MULTILINE)
    # List symbols (keep line breaks)
    text = re.sub(r"^\s*[-*•]\s*", "• ", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+\.\s*", "• ", text, flags=re.MULTILINE)
    # Remove bold/italic
    text = text.replace("**", "").replace("*", "")
    # Remove backticks
    text = text.replace("`", "'")
    # Clean up multiple line breaks (max 2)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Strip trailing spaces
    text = "\n".join(line.rstrip() for line in text.split("\n"))
    # Strip leading/trailing spaces
    text = text.strip()
    # Length limit with ...
    if len(text) > max_len:
        text = text[:max_len] + "…"
    return text


# Thread pool for CPU-intensive operations (transcription, etc.)
from concurrent.futures import ThreadPoolExecutor
audio_executor = ThreadPoolExecutor(max_workers=2, thread_name_prefix="audio_processing")


async def _action_indicator(context: ContextTypes.DEFAULT_TYPE, chat_id: int, action: str, stop_event: asyncio.Event):
    try:
        while not stop_event.is_set():
            try:
                await context.bot.send_chat_action(chat_id=chat_id, action=action)
            except Exception:
                pass
            # Telegram은 5초 동안 액션 유지. 4초 주기로 새로 송신.
            try:
                await asyncio.wait_for(stop_event.wait(), timeout=4.0)
            except asyncio.TimeoutError:
                continue
    except Exception:
        pass

class ActionIndicator:
    def __init__(self, context: ContextTypes.DEFAULT_TYPE, chat_id: int, action: str):
        self.context = context
        self.chat_id = chat_id
        self.action = action
        self.stop_event = asyncio.Event()
        self.task: asyncio.Task | None = None

    async def __aenter__(self):
        self.task = asyncio.create_task(_action_indicator(self.context, self.chat_id, self.action, self.stop_event))
        return self

    async def __aexit__(self, exc_type, exc, tb):
        self.stop_event.set()
        if self.task:
            try:
                await asyncio.wait_for(self.task, timeout=1.5)
            except Exception:
                self.task.cancel()


async def save_memory(user_id: str, username: str, message: str, response: str):
    if not supabase:
        return
    try:
        supabase.table("conversations").insert({
            "user_id": user_id,
            "username": username,
            "message": message,
            "response": response,
            "created_at": datetime.utcnow().isoformat()
        }).execute()
    except Exception as e:
        logger.warning(f"save_memory failed: {e}")


async def fetch_memory(user_id: str, limit: int = 8) -> List[Dict[str, str]]:
    if not supabase:
        return []
    try:
        res = supabase.table("conversations").select("message,response,created_at").eq("user_id", user_id).order("created_at", desc=True).limit(limit).execute()
        return list(reversed(res.data or []))
    except Exception as e:
        logger.warning(f"fetch_memory failed: {e}")
        return []


async def reply_text(update: Update, text: str):
    # Prevent telegram 409: retry with slight delay on 409
    try:
        await update.message.reply_text(text)
    except Exception as e:
        logger.warning(f"reply_text failed: {e}")
        await asyncio.sleep(0.8)
        try:
            await update.message.reply_text(text[:4000])
        except Exception:
            pass


async def handle_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Greet the user and surface primary capabilities."""
    return await text_handlers.handle_start(RUNTIME, update, context)


async def handle_mode(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /mode command interactions."""
    return await text_handlers.handle_mode(RUNTIME, update, context)


async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Main chat handler for free-form text conversation."""
    return await text_handlers.handle_text(RUNTIME, update, context)


async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle incoming photos with Gemini multimodal analysis."""
    return await media_handlers.handle_photo(RUNTIME, update, context)


async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle incoming voice messages with adaptive processing."""
    return await media_handlers.handle_voice(RUNTIME, update, context)


async def process_voice_background(update, context, chat_id, user_id, username, ack_msg):
    """Process voice in background - non-blocking, allows immediate responses."""
    return await media_handlers.process_voice_background(RUNTIME, update, context, chat_id, user_id, username, ack_msg)


async def process_with_gemini_multimodal(ogg_path: str, duration: float, chat_id: int, context, progress_messages):
    """Process short audio with Gemini 2.5 Flash multimodal."""
    return await media_handlers.process_with_gemini_multimodal(RUNTIME, ogg_path, duration, chat_id, context, progress_messages)


async def process_with_whisper_gemini(ogg_path: str, wav_path: str, duration: float, chat_id: int, context, progress_messages):
    """Process long audio with Whisper + Gemini."""
    return await media_handlers.process_with_whisper_gemini(RUNTIME, ogg_path, wav_path, duration, chat_id, context, progress_messages)


async def handle_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Display the user's recent document history."""
    return await text_handlers.handle_list(RUNTIME, update, context)


# ========== Google Drive Sync Handlers ==========

async def handle_drive(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /drive command - show Google Drive sync help."""
    return await drive_handlers.handle_drive(RUNTIME, update, context)


async def handle_drive_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /drivelist command - list all files in Google Drive."""
    return await drive_handlers.handle_drive_list(RUNTIME, update, context)


async def handle_drive_get(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /driveget command - download a file from Google Drive."""
    return await drive_handlers.handle_drive_get(RUNTIME, update, context)


async def handle_drive_sync(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /drivesync command - check for new files in Google Drive."""
    return await drive_handlers.handle_drive_sync(RUNTIME, update, context)

# ========== Gmail Handlers ==========


async def handle_gmail_on(update: Update, context: ContextTypes.DEFAULT_TYPE):
    return await gmail_handle_on(RUNTIME, update, context)


async def handle_gmail_off(update: Update, context: ContextTypes.DEFAULT_TYPE):
    return await gmail_handle_off(RUNTIME, update, context)


async def handle_gmail_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    return await gmail_handle_status(RUNTIME, update, context)


async def handle_gmail_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    return await gmail_handle_list(RUNTIME, update, context)

# ========== Gmail Monitoring Functions ==========

def start_gmail_monitoring():
    """Start Gmail monitoring in background thread"""
    import threading
    if gmail_monitoring_state["thread"] and gmail_monitoring_state["thread"].is_alive():
        return

    gmail_monitoring_state["thread"] = threading.Thread(
        target=gmail_monitor_loop,
        daemon=True
    )
    gmail_monitoring_state["thread"].start()
    logger.info("📧 Gmail monitoring started")


def gmail_monitor_loop():
    """Background Gmail monitoring loop"""
    import time

    try:
        from backend.services.gmail import GmailService
        gmail_service = GmailService()

        if not gmail_service.authenticate():
            logger.error("Gmail authentication failed")
            return

        logger.info("📧 Gmail monitoring worker started")

        while gmail_monitoring_state["enabled"]:
            try:
                logger.info("📧 Checking for new emails...")

                # Get recent emails
                recent_emails = gmail_service.get_recent_emails(max_results=20)
                new_emails = []

                for email_info in recent_emails:
                    email_id = email_info['id']

                    # Check if already processed
                    if email_id not in gmail_service.processed_emails:
                        email_content = gmail_service.get_email_content(email_id)
                        if email_content:
                            new_emails.append(email_content)
                            gmail_service.processed_emails.add(email_id)
                            try:
                                gmail_service.mark_as_read(email_id)
                            except Exception as mark_err:
                                logger.warning(f"Failed to mark email as read ({email_id}): {mark_err}")

                # Process new emails
                if new_emails:
                    logger.info(f"📧 Found {len(new_emails)} new emails")
                    gmail_monitoring_state["total_emails"] += len(new_emails)

                    for email_data in new_emails:
                        asyncio.run_coroutine_threadsafe(
                            process_and_send_email(email_data),
                            asyncio.get_event_loop()
                        )

                # Save processed emails
                gmail_service.save_processed_emails()
                gmail_monitoring_state["last_check"] = datetime.now().strftime("%H:%M:%S")

                # Wait 5 minutes
                for _ in range(300):  # Check every second for shutdown
                    if not gmail_monitoring_state["enabled"]:
                        break
                    time.sleep(1)

            except Exception as e:
                logger.error(f"Gmail monitoring error: {e}")
                time.sleep(60)  # Wait 1 minute on error

        logger.info("📧 Gmail monitoring worker stopped")

    except Exception as e:
        logger.error(f"Gmail monitoring loop error: {e}")


async def process_and_send_email(email_data):
    """Process email with Gemini and send to Telegram"""
    if not gemini_model:
        return

    try:
        # Create progress message
        progress_msg = await _app_instance.bot.send_message(
            chat_id=_app_instance.chat_ids[0] if _app_instance.chat_ids else None,
            text="📧 새 메일 분석 중..."
        )

        # Gemini summarization
        prompt = f"""
        다음 이메일을 한국어로 요약해주세요:

        보낸사람: {email_data['sender']}
        제목: {email_data['subject']}
        내용: {email_data['body']}

        요약 형식:
        - 핵심 내용 (2-3문장)
        - 중요도 (높음/보통/낮음)
        - 필요한 액션이 있다면 언급
        """

        response = gemini_model.generate_content(prompt)
        summary = format_plain(response.text)

        # Final message
        final_message = f"""
📧 **새 메일 요약**

👤 **보낸사람**: {email_data['sender']}
📝 **제목**: {email_data['subject']}
🕒 **시간**: {email_data['date']}

🤖 **AI 요약**:
{summary}
        """.strip()

        await _app_instance.bot.edit_message_text(
            chat_id=progress_msg.chat_id,
            message_id=progress_msg.message_id,
            text=final_message
        )

    except Exception as e:
        logger.error(f"Email processing error: {e}")


# ========== Calendar Handlers ==========

async def handle_cal_on(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_on command - Start Calendar monitoring."""
    return await calendar_handlers.handle_cal_on(RUNTIME, update, context)


async def handle_cal_off(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_off command - Stop Calendar monitoring."""
    return await calendar_handlers.handle_cal_off(RUNTIME, update, context)


async def handle_cal_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_status command - Check Calendar monitoring status."""
    return await calendar_handlers.handle_cal_status(RUNTIME, update, context)


async def handle_cal_today(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_today command - Show today's events."""
    return await calendar_handlers.handle_cal_today(RUNTIME, update, context)


async def handle_cal_tomorrow(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_tomorrow command - Show tomorrow's events."""
    return await calendar_handlers.handle_cal_tomorrow(RUNTIME, update, context)


async def handle_cal_week(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_week command - Show this week's events."""
    return await calendar_handlers.handle_cal_week(RUNTIME, update, context)


async def handle_cal_search(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Handle /cal_search command - Search for events."""
    return await calendar_handlers.handle_cal_search(RUNTIME, update, context)

# ========== Calendar Monitoring Functions ==========

def start_calendar_monitoring():
    """Start Calendar monitoring in background thread."""
    return calendar_handlers.start_calendar_monitoring(RUNTIME)


def calendar_monitor_loop():
    """Background Calendar monitoring loop."""
    return calendar_handlers.calendar_monitor_loop(RUNTIME)


async def process_and_send_calendar_alert(event_data):
    """Process event and send alert to Telegram."""
    return await calendar_handlers.process_and_send_calendar_alert(RUNTIME, event_data)

async def monitor_drive_changes():
    """Background task to monitor Google Drive for changes."""
    return await drive_handlers.monitor_drive_changes(RUNTIME)


def extract_text_from_file(file_path: str, file_name: str) -> str:
    """
    Extract text from various file formats
    Supports: txt, md, py, js, html, css, json, xml, csv, pdf, docx, pptx, xlsx, zip
    """
    import os
    import zipfile
    import chardet

    file_ext = os.path.splitext(file_name)[1].lower()

    try:
        # 1. Text-based files (most common)
        if file_ext in ['.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.html',
                        '.htm', '.css', '.scss', '.sass', '.less', '.json', '.xml',
                        '.csv', '.tsv', '.yaml', '.yml', '.ini', '.cfg', '.conf',
                        '.log', '.sql', '.sh', '.bat', '.ps1', '.dockerfile',
                        '.gitignore', '.env', '.properties', '.toml', '.r', '.R']:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                enc = chardet.detect(raw_data).get('encoding') or 'utf-8'
                return raw_data.decode(enc, errors='ignore')

        # 2. PDF files
        elif file_ext == '.pdf':
            try:
                import PyPDF2
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                return text
            except ImportError:
                return "PDF 읽기를 위한 PyPDF2가 설치되어 있지 않습니다."

        # 3. Word documents (.docx)
        elif file_ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                return text
            except ImportError:
                return "DOCX 읽기를 위한 python-docx가 설치되어 있지 않습니다."

        # 4. PowerPoint (.pptx)
        elif file_ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except ImportError:
                return "PPTX 읽기를 위한 python-pptx가 설치되어 있지 않습니다."

        # 5. Excel files (.xlsx, .xls)
        elif file_ext in ['.xlsx', '.xls']:
            try:
                import pandas as pd
                # Try to read all sheets
                df = pd.read_excel(file_path, sheet_name=None)
                text = ""
                for sheet_name, sheet_df in df.items():
                    text += f"\n=== Sheet: {sheet_name} ===\n"
                    text += sheet_df.to_string(index=False) + "\n"
                return text
            except ImportError:
                return "Excel 읽기를 위한 pandas가 설치되어 있지 않습니다."

        # 6. ZIP archives (extract and read text files inside)
        elif file_ext == '.zip':
            try:
                text = "=== ZIP 아카이브 내용 ===\n"
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    file_list = zip_ref.namelist()
                    text += f"총 {len(file_list)}개 파일\n\n"
                    for file_in_zip in file_list[:10]:  # Show first 10 files
                        text += f"• {file_in_zip}\n"
                        # If it's a text file, try to extract and read
                        if any(file_in_zip.lower().endswith(ext) for ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml']):
                            try:
                                content = zip_ref.read(file_in_zip).decode('utf-8', errors='ignore')
                                text += f"  내용 미리보기:\n{content[:500]}...\n"
                            except:
                                pass
                    if len(file_list) > 10:
                        text += f"\n... 외 {len(file_list) - 10}개 파일"
                return text
            except Exception as e:
                return f"ZIP 파일 읽기 오류: {str(e)}"

        # 7. Other binary files
        else:
            return f"지원하지 않는 파일 형식: {file_ext}\n파일 크기: {os.path.getsize(file_path)} bytes"

    except Exception as e:
        return f"파일 읽기 오류: {str(e)}"


async def handle_document_auto_save(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Auto-save incoming documents to Google Drive and analyze them."""
    return await drive_handlers.handle_document_auto_save(RUNTIME, update, context)


def main():
    print("=== 125 Unified Telegram Bot (Gemini 2.5 Flash + Drive Sync) ===")
    print(f"TELEGRAM_BOT_TOKEN: {'Set' if TELEGRAM_BOT_TOKEN else 'Not Found'}")
    print(f"GEMINI_API_KEY: {'Set' if GEMINI_API_KEY else 'Not Found'}")
    print(f"Supabase: {'Set' if (SUPABASE_URL and SUPABASE_KEY) else 'Not Set'}")
    print(f"Google Drive: {'Set' if os.path.exists(os.path.join(os.path.dirname(__file__), '..', '..', 'service_account.json')) else 'Not Set'}")
    print(f"Drive Monitoring: {'Enabled' if ENABLE_DRIVE_MONITORING else 'Disabled'} (interval: {DRIVE_MONITOR_INTERVAL}s)")

    if not TELEGRAM_BOT_TOKEN:
        print("ERROR: TELEGRAM_BOT_TOKEN is missing")
        return

    app = build_application(TELEGRAM_BOT_TOKEN)

    # Store app instance for Drive monitoring
    global _app_instance
    _app_instance = app

    # Allow new modular registrations to extend the application.
    from backend.bots import register_main_bot_handlers
    register_main_bot_handlers(app)

    logger.info("Handlers registered. Starting polling...")

    # Start Drive monitoring in a separate thread
    if ENABLE_DRIVE_MONITORING:
        def run_monitoring():
            import asyncio
            # Create new event loop for the thread
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            loop.run_until_complete(monitor_drive_changes())

        import threading
        monitor_thread = threading.Thread(target=run_monitoring, daemon=True)
        monitor_thread.start()
        logger.info("Drive monitoring worker started in background thread")

    app.run_polling()


if __name__ == "__main__":
    main()
