"""
125 Build Automation Enhanced - AI 문서 분석 봇 (확장 버전)
Gemini 기반 문서 분석 + RAG + 텔레그램/드라이브 통합
"""
import os
import asyncio
import logging
from datetime import datetime
from typing import Optional, Dict, List, Any

# 외부 라이브러리
import google.generativeai as genai
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
import telegram
from telegram import Update, File
from telegram.ext import Application, CommandHandler, MessageHandler, ContextTypes, filters
import httpx
import chardet
import markdown_it
from bs4 import BeautifulSoup
import csv
import openpyxl
from pptx import Presentation
import json
import hashlib
from pathlib import Path

# 환경변수 로드
from dotenv import load_dotenv
load_dotenv()

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('backend.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# 상수 정의
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
TELEGRAM_BOT_TOKEN = os.getenv('TELEGRAM_BOT_TOKEN')
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
ENABLE_RAG = os.getenv('ENABLE_RAG', 'false').lower() == 'true'
VECTOR_STORE_PATH = os.getenv('VECTOR_STORE_PATH', 'data/store')
GEN_TEMPERATURE = float(os.getenv('GEN_TEMPERATURE', '0.2'))
GEN_MAX_OUTPUT_TOKENS = int(os.getenv('GEN_MAX_OUTPUT_TOKENS', '2048'))

# 조건부 import (RAG 활성화 시에만)
faiss = None
chromadb = None
SentenceTransformer = None

if ENABLE_RAG:
    try:
        import faiss
        import chromadb
        from sentence_transformers import SentenceTransformer
    except ImportError as e:
        logger.warning(f"RAG 관련 모듈 import 실패: {e}")
        ENABLE_RAG = False

# Gemini 모델 초기화 (API 키가 있을 때만)
model = None
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    generation_config = genai.GenerationConfig(
        temperature=GEN_TEMPERATURE,
        top_p=0.9,
        max_output_tokens=GEN_MAX_OUTPUT_TOKENS
    )
    model = genai.GenerativeModel('gemini-2.5-flash', generation_config=generation_config)

# 글로벌 변수
drive_service = None
telegram_app = None
vector_store = None
embedding_model = None

async def init_services():
    """서비스 초기화"""
    global drive_service, telegram_app, vector_store, embedding_model

    # Gemini 확인
    if not GEMINI_API_KEY:
        logger.error("GEMINI_API_KEY가 설정되지 않았습니다")
        return False

    # RAG 초기화 (선택)
    if ENABLE_RAG:
        try:
            embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
            Path(VECTOR_STORE_PATH).mkdir(parents=True, exist_ok=True)
            vector_store = chromadb.PersistentClient(path=VECTOR_STORE_PATH)
            logger.info("RAG 시스템 초기화 완료")
        except Exception as e:
            logger.error(f"RAG 초기화 실패: {e}")

    logger.info("모든 서비스 초기화 완료")
    return True

# ===== 범용 문서 추출기 =====

def extract_text_from_markdown(path: str) -> str:
    """Markdown 파일에서 텍스트 추출"""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()

        # markdown-it으로 HTML 변환 후 BeautifulSoup으로 텍스트 추출
        md = markdown_it.MarkdownIt()
        html = md.render(content)
        soup = BeautifulSoup(html, 'html.parser')

        # 헤딩, 목록 등 구조 유지하면서 텍스트 추출
        return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logger.warning(f"Markdown 추출 실패, UTF-8 재시도: {e}")
        try:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e2:
            logger.error(f"Markdown 추출 최종 실패: {e2}")
            return ""

def extract_text_from_html(path: str) -> str:
    """HTML 파일에서 텍스트 추출"""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()

        soup = BeautifulSoup(content, 'html.parser')

        # 스크립트, 스타일 제거
        for script in soup(["script", "style"]):
            script.decompose()

        # 본문 텍스트만 추출
        return soup.get_text(separator='\n', strip=True)
    except Exception as e:
        logger.error(f"HTML 추출 실패: {e}")
        return ""

def extract_text_from_csv(path: str) -> str:
    """CSV 파일에서 텍스트 추출"""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            rows = list(reader)

        # 헤더와 데이터 결합
        text_parts = []
        if rows:
            text_parts.append("CSV Headers: " + ", ".join(rows[0]))
            for i, row in enumerate(rows[1:], 1):
                text_parts.append(f"Row {i}: " + ", ".join(row))

        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"CSV 추출 실패: {e}")
        return ""

def extract_text_from_xlsx(path: str) -> str:
    """Excel 파일에서 텍스트 추출"""
    try:
        wb = openpyxl.load_workbook(path, read_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")

            for row in sheet.iter_rows(values_only=True):
                # None이 아닌 값만 필터링
                row_data = [str(cell) for cell in row if cell is not None]
                if row_data:
                    text_parts.append(", ".join(row_data))

        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"XLSX 추출 실패: {e}")
        return ""

def extract_text_from_pptx(path: str) -> str:
    """PowerPoint 파일에서 텍스트 추출"""
    try:
        prs = Presentation(path)
        text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {i}:")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)

        return "\n".join(text_parts)
    except Exception as e:
        logger.error(f"PPTX 추출 실패: {e}")
        return ""

def extract_text_from_pdf(path: str) -> str:
    """PDF 파일에서 텍스트 추출"""
    try:
        import PyPDF2
        with open(path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        # pdfplumber로 재시도
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                text = ""
                for page in pdf.pages:
                    text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e2:
            logger.error(f"PDF 추출 실패: {e2}")
            return ""

def extract_text_fallback(path: str) -> str:
    """텍스트 파일 추출 (chardet로 인코딩 추정)"""
    try:
        with open(path, 'rb') as f:
            raw_data = f.read()

        # 인코딩 추정
        detected = chardet.detect(raw_data)
        encoding = detected.get('encoding', 'utf-8')

        # 추정된 인코딩으로 디코딩
        return raw_data.decode(encoding, errors='ignore')
    except Exception as e:
        logger.error(f"Fallback 추출 실패: {e}")
        return ""

def get_text_extractor(mime_type: str, file_path: str) -> str:
    """MIME 타입에 따른 텍스트 추출기 선택"""
    mime_to_extractor = {
        'text/markdown': extract_text_from_markdown,
        'text/html': extract_text_from_html,
        'text/csv': extract_text_from_csv,
        'application/pdf': extract_text_from_pdf,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': extract_text_from_xlsx,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_text_from_pptx,
    }

    # 확장자 기반 추가 매핑
    ext_to_extractor = {
        '.md': extract_text_from_markdown,
        '.markdown': extract_text_from_markdown,
        '.html': extract_text_from_html,
        '.htm': extract_text_from_html,
        '.csv': extract_text_from_csv,
        '.pdf': extract_text_from_pdf,
        '.xlsx': extract_text_from_xlsx,
        '.pptx': extract_text_from_pptx,
    }

    # MIME 타입 우선
    if mime_type in mime_to_extractor:
        return mime_to_extractor[mime_type](file_path)

    # 확장자 기반
    ext = Path(file_path).suffix.lower()
    if ext in ext_to_extractor:
        return ext_to_extractor[ext](file_path)

    # 기타 텍스트 파일
    if mime_type.startswith('text/') or ext in ['.txt', '.log', '.json', '.xml']:
        return extract_text_fallback(file_path)

    return ""

# ===== 청크 + 맵리듀스 요약 유틸리티 =====

def split_into_chunks(text: str, chunk_chars: int = 4000, overlap: int = 400) -> List[str]:
    """텍스트를 겹치는 청크로 분할"""
    if len(text) <= chunk_chars:
        return [text]

    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_chars

        # 단어 경계에서 자르기
        if end < len(text):
            # 공백이나 줄바꿈에서 자르기
            while end > start and text[end] not in [' ', '\n', '\t']:
                end -= 1
            if end == start:  # 단어 경계 못 찾음
                end = start + chunk_chars

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - overlap
        if start >= len(text):
            break

    return chunks

def summarize_chunk(chunk: str) -> str:
    """단일 청크 요약"""
    prompt = f"""역할: 전문가 보조 에이전트

다음 텍스트를 분석하여 핵심 내용을 요약해주세요.

요약 지침:
- 섹션별로 구조화: 요약/핵심포인트/액션아이템/날짜/리스크
- 근거가 약하면 '추정'으로 표기
- 간결하고 구조화된 형식으로 작성

텍스트:
{chunk}

요약:"""

    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        logger.error(f"청크 요약 실패: {e}")
        return f"요약 실패: {chunk[:200]}..."

def compose_summaries(summaries: List[str]) -> str:
    """청크 요약들을 통합 요약"""
    combined = "\n\n".join(f"청크 {i+1}: {summary}" for i, summary in enumerate(summaries))

    prompt = f"""역할: 전문가 보조 에이전트

다음은 여러 청크의 요약입니다. 이를 종합하여 전체 문서의 통합 요약을 작성해주세요.

통합 요약 지침:
- 전체 문서의 주요 테마와 내용을 포괄
- 섹션별 구조화 유지
- 중복 제거 및 일관성 확보
- 핵심 인사이트 강조

청크 요약들:
{combined}

통합 요약:"""

    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        logger.error(f"통합 요약 실패: {e}")
        return "통합 요약 실패: " + " ".join(summaries)

def summarize_text(text: str) -> str:
    """텍스트를 청크로 분할하여 요약 (필요시)"""
    if len(text) <= 4000:
        # 단일 패스 요약
        return summarize_chunk(text)

    # 청크 + 맵리듀스
    chunks = split_into_chunks(text)
    summaries = [summarize_chunk(chunk) for chunk in chunks]
    return compose_summaries(summaries)

# ===== Google Drive 통합 =====

async def init_drive_service():
    """Google Drive API 서비스 초기화"""
    global drive_service

    try:
        creds = None
        token_path = 'token.json'

        if os.path.exists(token_path):
            creds = Credentials.from_authorized_user_file(token_path, SCOPES)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(
                    'credentials.json', SCOPES)
                creds = flow.run_local_server(port=0)

            with open(token_path, 'w') as token:
                token.write(creds.to_json())

        drive_service = build('drive', 'v3', credentials=creds)
        logger.info("Google Drive 서비스 초기화 완료")
        return True

    except Exception as e:
        logger.error(f"Drive 서비스 초기화 실패: {e}")
        return False

def analyze_drive_file(file_id: str, mime_type: str, file_name: str) -> Dict[str, Any]:
    """Drive 파일 분석 및 요약"""
    try:
        # 파일 다운로드
        request = drive_service.files().get_media(fileId=file_id)
        file_path = f"/tmp/{file_id}_{file_name}"

        with open(file_path, 'wb') as f:
            f.write(request.execute())

        # 텍스트 추출
        text = get_text_extractor(mime_type, file_path)

        if not text:
            return {
                'success': False,
                'error': '텍스트 추출 실패',
                'file_name': file_name
            }

        # 요약
        summary = summarize_text(text)

        # 임시 파일 정리
        os.remove(file_path)

        return {
            'success': True,
            'file_name': file_name,
            'summary': summary,
            'text_length': len(text)
        }

    except Exception as e:
        logger.error(f"Drive 파일 분석 실패: {e}")
        return {
            'success': False,
            'error': str(e),
            'file_name': file_name
        }

# ===== Telegram 봇 핸들러 =====

# 글로벌 변수에 최근 문서 저장
recent_documents = {}

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """텔레그램 문서 핸들러 - 문서 저장만 수행"""
    try:
        document = update.message.document
        if not document:
            return

        file_name = document.file_name
        mime_type = document.mime_type
        file_id = document.file_id
        user_id = update.effective_user.id

        # 지원하지 않는 형식 체크
        supported_extensions = ['.md', '.markdown', '.html', '.csv', '.pdf', '.xlsx', '.pptx', '.txt', '.log', '.json']
        file_ext = Path(file_name).suffix.lower()

        if file_ext not in supported_extensions and not mime_type.startswith('text/'):
            await update.message.reply_text(
                f"❌ 지원하지 않는 파일 형식입니다: {file_ext}\n"
                "지원 형식: .md, .html, .csv, .pdf, .xlsx, .pptx, .txt, .log, .json"
            )
            return

        # 파일 다운로드 및 임시 저장
        file = await context.bot.get_file(file_id)
        file_path = f"/tmp/{file_id}_{file_name}"

        await file.download_to_drive(file_path)

        # 사용자별 최근 문서 저장
        if user_id not in recent_documents:
            recent_documents[user_id] = []

        # 텍스트 추출
        text = get_text_extractor(mime_type, file_path)

        if not text:
            await update.message.reply_text("❌ 텍스트 추출 실패 (파일 형식이 지원되지 않거나 손상됨)")
            os.remove(file_path)
            return

        doc_info = {
            'file_name': file_name,
            'file_path': file_path,
            'mime_type': mime_type,
            'text': text,
            'text_length': len(text),
            'timestamp': datetime.now()
        }

        recent_documents[user_id].append(doc_info)

        # RAG 저장 (활성화된 경우)
        if ENABLE_RAG:
            await rag_store_document(file_path, file_name, text, str(user_id))

        # 최대 5개까지만 저장
        if len(recent_documents[user_id]) > 5:
            old_doc = recent_documents[user_id].pop(0)
            if os.path.exists(old_doc['file_path']):
                os.remove(old_doc['file_path'])

        await update.message.reply_text(
            f"- **문서 저장 완료**\n\n"
            f"- **파일명:** {file_name}\n"
            f"- **형식:** {mime_type}\n\n"
            f"- 분석을 원하시면 다음 명령을 사용하세요:\n"
            f"- `/analyze` - 최근 문서 분석\n"
            f"- `/summarize` - 최근 문서 요약\n"
            f"- `/ask [질문]` - 문서에 대해 질문",
            parse_mode='Markdown'
        )

    except Exception as e:
        logger.error(f"문서 저장 실패: {e}")
        await update.message.reply_text("❌ 문서 저장 중 오류가 발생했습니다")

async def handle_analyze(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """문서 분석 핸들러 (/analyze 명령)"""
    user_id = update.effective_user.id

    if user_id not in recent_documents or not recent_documents[user_id]:
        await update.message.reply_text("❌ 최근에 업로드한 문서가 없습니다. 먼저 문서를 업로드해주세요.")
        return

    try:
        # 가장 최근 문서 분석
        latest_doc = recent_documents[user_id][-1]

        await update.message.reply_text("🔍 문서를 분석하고 있습니다...")

        # AI 분석 수행
        analysis_prompt = f"""역할: 전문 문서 분석가

다음 문서를 전문적으로 분석해주세요.

분석 요구사항:
- 문서의 주요 목적과 내용 파악
- 구조와 구성 요소 분석
- 핵심 개념과 주요 포인트 도출
- 잠재적 활용 방안 제시
- 개선점이나 주의사항 언급

문서 정보:
- 파일명: {latest_doc['file_name']}
- 형식: {latest_doc['mime_type']}
- 길이: {latest_doc['text_length']}자

문서 내용:
{latest_doc['text']}

분석 결과:"""

        response = model.generate_content(analysis_prompt)
        analysis = response.text.strip()

        response_msg = f"📊 **문서 분석 결과**\n\n**파일:** {latest_doc['file_name']}\n\n{analysis}"

        if len(response_msg) > 4000:
            response_msg = response_msg[:3997] + "..."

        await update.message.reply_text(response_msg, parse_mode='Markdown')

    except Exception as e:
        logger.error(f"문서 분석 실패: {e}")
        await update.message.reply_text("❌ 문서 분석 중 오류가 발생했습니다")

async def handle_summarize(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """문서 요약 핸들러 (/summarize 명령)"""
    user_id = update.effective_user.id

    if user_id not in recent_documents or not recent_documents[user_id]:
        await update.message.reply_text("❌ 최근에 업로드한 문서가 없습니다. 먼저 문서를 업로드해주세요.")
        return

    try:
        # 가장 최근 문서 요약
        latest_doc = recent_documents[user_id][-1]

        await update.message.reply_text("📝 문서를 요약하고 있습니다...")

        # 요약 수행
        summary = summarize_text(latest_doc['text'])

        response_msg = f"📄 **문서 요약 결과**\n\n**파일:** {latest_doc['file_name']}\n**텍스트 길이:** {latest_doc['text_length']}자\n\n{summary}"

        if len(response_msg) > 4000:
            response_msg = response_msg[:3997] + "..."

        await update.message.reply_text(response_msg, parse_mode='Markdown')

    except Exception as e:
        logger.error(f"문서 요약 실패: {e}")
        await update.message.reply_text("❌ 문서 요약 중 오류가 발생했습니다")

async def handle_list(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """저장된 문서 목록 핸들러 (/list 명령)"""
    user_id = update.effective_user.id

    if user_id not in recent_documents or not recent_documents[user_id]:
        await update.message.reply_text("📂 저장된 문서가 없습니다.")
        return

    try:
        doc_list = []
        for i, doc in enumerate(recent_documents[user_id], 1):
            timestamp = doc['timestamp'].strftime('%H:%M:%S')
            doc_list.append(f"{i}. {doc['file_name']} ({doc['text_length']}자) - {timestamp}")

        response = "📂 **저장된 문서 목록**\n\n" + "\n".join(doc_list)
        response += f"\n\n총 {len(recent_documents[user_id])}개 문서가 저장되어 있습니다."

        await update.message.reply_text(response, parse_mode='Markdown')

    except Exception as e:
        logger.error(f"문서 목록 조회 실패: {e}")
        await update.message.reply_text("❌ 문서 목록 조회 중 오류가 발생했습니다")

async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """텍스트 메시지 핸들러"""
    user_id = update.effective_user.id
    message_text = update.message.text.strip()

    # 명령어 인식
    if message_text.startswith('/'):
        return  # 명령어는 별도 핸들러에서 처리

    # 일반 질문 처리
    if user_id in recent_documents and recent_documents[user_id]:
        try:
            # RAG가 활성화된 경우
            if ENABLE_RAG and vector_store:
                answer = await rag_query(message_text, str(user_id))
                response = f"🤖 **질문:** {message_text}\n\n**답변:**\n{answer}"
            else:
                # 최근 문서에 대해 일반 AI 질문
                latest_doc = recent_documents[user_id][-1]

                prompt = f"""사용자의 질문에 대해 최근 업로드된 문서의 내용을 참고하여 답변해주세요.

문서 정보:
- 파일명: {latest_doc['file_name']}
- 내용: {latest_doc['text'][:2000]}... (축약)

질문: {message_text}

답변:"""

                response = model.generate_content(prompt)
                answer = response.text.strip()
                response = f"🤖 **질문:** {message_text}\n\n**답변:**\n{answer}"

            if len(response) > 4000:
                response = response[:3997] + "..."

            await update.message.reply_text(response, parse_mode='Markdown')

        except Exception as e:
            logger.error(f"텍스트 질문 처리 실패: {e}")
            await update.message.reply_text("❌ 질문 처리 중 오류가 발생했습니다")
    else:
        # 기본 도움말
        help_msg = """🤖 **125 Build Automation 봇**

사용 가능한 기능:
- 문서 업로드 후 다음 명령어 사용:
  - `/analyze` - 문서 전문 분석
  - `/summarize` - 문서 요약
  - `/list` - 저장된 문서 목록
  - `/ask [질문]` - RAG 기반 질문

- 또는 일반 텍스트로 질문하기

먼저 문서를 업로드해주세요!"""
        await update.message.reply_text(help_msg, parse_mode='Markdown')

async def handle_ask(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """질문 핸들러 (/ask 명령)"""
    if not ENABLE_RAG or not vector_store:
        await update.message.reply_text("❌ RAG 시스템이 활성화되지 않았습니다")
        return

    try:
        query = " ".join(context.args)
        if not query:
            await update.message.reply_text("❌ 질문을 입력해주세요: /ask [질문]")
            return

        user_id = update.effective_user.id
        # RAG 검색 및 답변 생성
        answer = await rag_query(query, str(user_id))

        await update.message.reply_text(f"🤖 **질문:** {query}\n\n**답변:**\n{answer}", parse_mode='Markdown')

    except Exception as e:
        logger.error(f"질문 처리 실패: {e}")
        await update.message.reply_text("❌ 질문 처리 중 오류가 발생했습니다")

# ===== RAG 시스템 =====

async def rag_store_document(file_path: str, file_name: str, text: str, owner_id: str):
    """문서를 벡터 스토어에 저장"""
    if not ENABLE_RAG or not vector_store or not embedding_model:
        return False

    try:
        collection_name = f"docs_{owner_id}"
        collection = vector_store.get_or_create_collection(name=collection_name)

        # 청크 분할
        chunks = split_into_chunks(text, chunk_chars=1000, overlap=100)

        # 임베딩 및 저장
        for i, chunk in enumerate(chunks):
            chunk_id = f"{file_name}_{i}_{hashlib.md5(chunk.encode()).hexdigest()[:8]}"
            embedding = embedding_model.encode(chunk).tolist()

            metadata = {
                'file_name': file_name,
                'chunk_index': i,
                'owner_id': owner_id,
                'created_at': datetime.now().isoformat()
            }

            collection.add(
                ids=[chunk_id],
                embeddings=[embedding],
                metadatas=[metadata],
                documents=[chunk]
            )

        logger.info(f"RAG 저장 완료: {file_name} ({len(chunks)} 청크)")
        return True

    except Exception as e:
        logger.error(f"RAG 저장 실패: {e}")
        return False

async def rag_query(query: str, owner_id: str = None, top_k: int = 3) -> str:
    """RAG 쿼리 수행"""
    if not ENABLE_RAG or not vector_store or not embedding_model:
        return "RAG 시스템이 비활성화되었습니다"

    try:
        collection_name = f"docs_{owner_id}" if owner_id else "docs_default"
        collection = vector_store.get_or_create_collection(name=collection_name)

        # 쿼리 임베딩
        query_embedding = embedding_model.encode(query).tolist()

        # 유사도 검색
        results = collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )

        if not results['documents']:
            return "관련 문서를 찾을 수 없습니다"

        # 컨텍스트 구성
        context_parts = []
        for i, (doc, metadata) in enumerate(zip(results['documents'][0], results['metadatas'][0])):
            context_parts.append(f"문서 {i+1} ({metadata['file_name']}):\n{doc}")

        context = "\n\n".join(context_parts)

        # 답변 생성
        prompt = f"""문서 근거 인용(파일명/슬라이드/시트/페이지/섹션 헤딩)
근거가 없으면 '모르겠다'로 응답
간결하게, 목록 위주

질문: {query}

참고 문서:
{context}

답변:"""

        response = model.generate_content(prompt)
        return response.text.strip()

    except Exception as e:
        logger.error(f"RAG 쿼리 실패: {e}")
        return f"RAG 쿼리 오류: {str(e)}"

# ===== 텔레그램 봇 전용 실행 함수 =====

async def run_telegram_bot():
    """텔레그램 봇만 실행"""
    logger.info("텔레그램 봇 시작")

    # 서비스 초기화
    if not await init_services():
        logger.error("서비스 초기화 실패")
        return

    # Drive 서비스 초기화 (선택)
    if os.path.exists('credentials.json'):
        await init_drive_service()

    # 텔레그램 봇 초기화
    if TELEGRAM_BOT_TOKEN:
        telegram_app = Application.builder().token(TELEGRAM_BOT_TOKEN).build()

        # 핸들러 등록
        telegram_app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
        telegram_app.add_handler(CommandHandler("analyze", handle_analyze))
        telegram_app.add_handler(CommandHandler("summarize", handle_summarize))
        telegram_app.add_handler(CommandHandler("list", handle_list))
        telegram_app.add_handler(CommandHandler("ask", handle_ask))
        telegram_app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

        logger.info("텔레그램 봇 핸들러 등록 완료")

        # 봇 실행
        await telegram_app.run_polling()
    else:
        logger.warning("TELEGRAM_BOT_TOKEN이 설정되지 않아 텔레그램 봇을 시작하지 않습니다")

    logger.info("텔레그램 봇 종료")

# ===== 메인 함수 =====

async def main():
    """메인 실행 함수"""
    logger.info("125 Build Automation Enhanced 시작")

    # 서비스 초기화
    if not await init_services():
        logger.error("서비스 초기화 실패")
        return

    # Drive 서비스 초기화 (선택)
    if os.path.exists('credentials.json'):
        await init_drive_service()

    # 텔레그램 봇 초기화
    if TELEGRAM_BOT_TOKEN:
        telegram_app = Application.builder().token(TELEGRAM_BOT_TOKEN).build()

        # 핸들러 등록
        telegram_app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
        telegram_app.add_handler(CommandHandler("analyze", handle_analyze))
        telegram_app.add_handler(CommandHandler("summarize", handle_summarize))
        telegram_app.add_handler(CommandHandler("list", handle_list))
        telegram_app.add_handler(CommandHandler("ask", handle_ask))
        telegram_app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

        logger.info("텔레그램 봇 핸들러 등록 완료")

        # 봇 실행
        await telegram_app.run_polling()
    else:
        logger.warning("TELEGRAM_BOT_TOKEN이 설정되지 않아 텔레그램 봇을 시작하지 않습니다")

    logger.info("125 Build Automation Enhanced 종료")

if __name__ == "__main__":
    print("125 Build Automation Enhanced - Fixed Bot (Final)")
    try:
        # Create a proper event loop with cleanup
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        
        try:
            loop.run_until_complete(run_telegram_bot())
        except KeyboardInterrupt:
            print("\nINFO: Bot stopped by user")
        except Exception as e:
            print(f"CRITICAL ERROR: {e}")
            import traceback
            traceback.print_exc()
        finally:
            # Clean up pending tasks
            pending = asyncio.all_tasks(loop)
            for task in pending:
                task.cancel()
            
            if pending:
                try:
                    loop.run_until_complete(asyncio.gather(*pending, return_exceptions=True))
                except Exception:
                    pass
            
            loop.close()
            
    except Exception as e:
        print(f"FATAL ERROR: {e}")
        sys.exit(1)
